#!/usr/bin/env python3
import sys
import os
import subprocess
import time
import json
import re
import mimetypes
import datetime
import hashlib
import pickle
import argparse
from urllib.parse import unquote
from typing import List, Optional

# External dependencies may fail before bootstrap ensures they are in the venv
try:
    import requests
    import numpy as np
    from adelaide_bridge import AdelaideBridge
except ImportError:
    requests = None
    np = None
    AdelaideBridge = None

try:
    import fitz # PyMuPDF
except ImportError:
    fitz = None

# --- Environment Setup ---
def apply_base_env():
    """Load core environment variables from config.json to ensure consistent execution."""
    config_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "config.json")
    if os.path.exists(config_path):
        try:
            with open(config_path, 'r') as f:
                config = json.load(f)
                base_env = config.get("base_env", {})
                for key, value in base_env.items():
                    os.environ[key] = value
        except Exception as e:
            print(f"⚠️ Error loading base_env: {e}", file=sys.stderr)

# --- Bootstrap Virtual Environment ---
VENV_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "pyvenv")
REQUIREMENTS = [
    "requests", "numpy", "Pillow", "PyMuPDF",
    "openpyxl", "python-docx", "python-pptx", "tinytag"
]

def bootstrap_venv():
    """Ensures the script runs in its dedicated virtual environment."""
    apply_base_env()
    venv_abs = os.path.abspath(VENV_DIR)
    
    if os.path.abspath(sys.prefix) != venv_abs:
        if not os.path.exists(VENV_DIR):
            print(f"[*] Creating virtual environment in {VENV_DIR}...", file=sys.stderr)
            subprocess.run([sys.executable, "-m", "venv", VENV_DIR], check=True)
            
        python_exe = os.path.join(VENV_DIR, "bin", "python") if os.name != 'nt' else os.path.join(VENV_DIR, "Scripts", "python.exe")
        
        if os.path.exists(python_exe):
            os.execv(python_exe, [python_exe] + sys.argv)

    import importlib.util
    # Note: Pillow is imported as PIL, PyMuPDF as fitz, python-docx as docx, python-pptx as pptx
    CHECK_MODULES = ["requests", "numpy", "PIL", "fitz", "openpyxl", "docx", "pptx", "tinytag"]
    missing = [mod for mod in CHECK_MODULES if importlib.util.find_spec(mod) is None]
    
    if missing:
        print(f"[*] Missing dependencies. Installing: {', '.join(REQUIREMENTS)}...", file=sys.stderr)
        pip_exe = os.path.join(VENV_DIR, "bin", "pip") if os.name != 'nt' else os.path.join(VENV_DIR, "Scripts", "pip.exe")
        subprocess.run([pip_exe, "install", "--upgrade", "pip"], check=True)
        subprocess.run([pip_exe, "install"] + REQUIREMENTS, check=True)
        os.execv(sys.executable, [sys.executable] + sys.argv)

bootstrap_venv()

# --- Post-Bootstrap Environment Fixes ---
if "RECOLL_CONFDIR" not in os.environ:
    os.environ["RECOLL_CONFDIR"] = os.path.expanduser("~/.recoll")

# ================= CONFIGURATION =================
recoll_cmd = "/Applications/Recoll.app/Contents/MacOS/recollq"
OLLAMA_BASE_URL = os.environ.get("OLLAMA_PROXY_URL", "http://localhost:11435")
OLLAMA_EMBED_ENDPOINT = f"{OLLAMA_BASE_URL}/api/embed"
OLLAMA_MODEL = "qwen3-embedding:0.6b"

TOP_FILES_TO_PROCESS = 10        # Aggressive Lexical Cutoff
TOP_CHUNKS_TO_RETURN = 5         # Final chunks to present
RANK_THRESHOLD = 0.55            # Minimum cosine similarity to display result
CHUNK_SIZE = 512
CHUNK_OVERLAP = 50
MAX_CHARS_PER_FILE = 100000

CACHE_FILE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "embed_cache.pkl")
MAX_CACHE_ENTRIES = 120000       # Mathematically approximates 512 MiB limit
# =================================================

# --- MEMORY CACHE LOGIC ---
MEMORY_CACHE = {}
CACHE_MODIFIED = False

def load_cache():
    global MEMORY_CACHE
    if os.path.exists(CACHE_FILE_PATH):
        try:
            with open(CACHE_FILE_PATH, 'rb') as f:
                MEMORY_CACHE = pickle.load(f)
            print(f"[*] Loaded {len(MEMORY_CACHE)} embedding vectors into active memory.", file=sys.stderr)
        except Exception as e:
            print(f"⚠️ Failed to load memory cache. Starting fresh: {e}", file=sys.stderr)
            MEMORY_CACHE = {}

def save_cache():
    global MEMORY_CACHE, CACHE_MODIFIED
    if not CACHE_MODIFIED:
        return
        
    if len(MEMORY_CACHE) > MAX_CACHE_ENTRIES:
        print(f"[*] Memory cache exceeded {MAX_CACHE_ENTRIES} entries. Executing LRU eviction...", file=sys.stderr)
        sorted_keys = sorted(MEMORY_CACHE.keys(), key=lambda k: MEMORY_CACHE[k]['last_used'])
        keys_to_delete = sorted_keys[:int(MAX_CACHE_ENTRIES * 0.2)]
        for k in keys_to_delete:
            del MEMORY_CACHE[k]
            
    try:
        with open(CACHE_FILE_PATH, 'wb') as f:
            pickle.dump(MEMORY_CACHE, f)
        print("[*] Memory cache flushed to disk successfully.", file=sys.stderr)
    except Exception as e:
        print(f"⚠️ Failed to write cache to disk: {e}", file=sys.stderr)

def get_embedding(text: str) -> Optional[np.ndarray]:
    global MEMORY_CACHE, CACHE_MODIFIED
    if not text or not text.strip():
        return None
    
    text_hash = hashlib.sha256(text.encode('utf-8')).hexdigest()
    
    if text_hash in MEMORY_CACHE:
        MEMORY_CACHE[text_hash]['last_used'] = time.time()
        return MEMORY_CACHE[text_hash]['embedding']

    try:
        resp = requests.post(OLLAMA_EMBED_ENDPOINT, json={"model": OLLAMA_MODEL, "input": text}, timeout=30)
        resp.raise_for_status()
        data = resp.json()
        
        if "embeddings" in data and len(data["embeddings"]) > 0:
            vector = data["embeddings"][0]
        elif "embedding" in data:
            vector = data["embedding"]
        else:
            raise KeyError("Neither 'embedding' nor 'embeddings' found in response")

        emb_array = np.array(vector, dtype=np.float32)
        
        MEMORY_CACHE[text_hash] = {
            'embedding': emb_array,
            'last_used': time.time()
        }
        CACHE_MODIFIED = True
        return emb_array
    except Exception as e:
        print(f"⚠️ Embedding API failed: {e}", file=sys.stderr)
        return None

# --- MAIN LOGIC ---
def ensure_ollama_running():
    try:
        requests.get(f"{OLLAMA_BASE_URL}", timeout=2)
        return True
    except Exception:
        print(f"⚠️ Proxy Ollama not reachable at {OLLAMA_BASE_URL}. Assuming it's managed externally or down.", file=sys.stderr)
        return False

def cosine_similarity(v1: np.ndarray, v2: np.ndarray) -> float:
    if v1 is None or v2 is None:
        return 0.0
    try:
        if AdelaideBridge:
            bridge = AdelaideBridge.get_instance()
            sim = bridge.cosine_similarity(v1, v2)
            if sim is not None:
                return sim
    except Exception:
        pass

    norm = (np.linalg.norm(v1) * np.linalg.norm(v2))
    return np.dot(v1, v2) / norm if norm != 0 else 0.0

def get_file_paths_from_massive_dump(query: str, limit: int) -> List[str]:
    cmd = [recoll_cmd, "-o", query, "-A", "-m", "-C", "-P", "-d"]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        pattern = re.compile(r'\[file://(.*?)\]')
        matches = pattern.findall(result.stdout)
        
        # Preserve Recoll's native ranking order while deduplicating
        unique_paths = []
        seen = set()
        for m in matches:
            decoded_path = unquote(m)
            if decoded_path not in seen:
                seen.add(decoded_path)
                unique_paths.append(decoded_path)
                if len(unique_paths) >= limit:
                    break
                    
        return unique_paths
    except subprocess.CalledProcessError as e:
        print(f"❌ recollq failed: {e.stderr}", file=sys.stderr)
        sys.exit(e.returncode)

def extract_content_via_python(path: str) -> str:
    if not os.path.exists(path):
        return ""
    ext = os.path.splitext(path)[1].lower()
    text = ""
    print(f"   ↳ Processing natively: {ext or 'Unknown/Text'}", file=sys.stderr)

    try:
        if ext == '.pdf' and fitz:
            doc = fitz.open(path)
            for page in doc:
                text += page.get_text() + "\n"
        elif ext in ['.xlsx', '.xls']:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            for sheet in wb.worksheets:
                text += f"\n--- Sheet: {sheet.title} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) for cell in row if cell is not None]
                    if row_data:
                        text += " | ".join(row_data) + "\n"
        elif ext in ['.docx']:
            import docx
            doc = docx.Document(path)
            text = "\n".join([p.text for p in doc.paragraphs if p.text])
        elif ext in ['.pptx']:
            import pptx
            prs = pptx.Presentation(path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        elif ext in ['.mp3', '.wav', '.mp4', '.mkv', '.flac', '.m4a']:
            from tinytag import TinyTag
            tag = TinyTag.get(path)
            text = f"[Media Metadata]\nFile: {os.path.basename(path)}\n"
            if tag.title:
                text += f"Title: {tag.title}\n"
            if tag.artist:
                text += f"Artist: {tag.artist}\n"
            if tag.album:
                text += f"Album: {tag.album}\n"
            if tag.year:
                text += f"Year: {tag.year}\n"
            if tag.duration:
                text += f"Duration: {int(tag.duration)} seconds\n"
        else:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
    except Exception as e:
        print(f"   ⚠️ Native extraction failed for {os.path.basename(path)}: {e}", file=sys.stderr)
    return text

def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    chunks = []
    if len(text) <= size:
        return [text]
    for i in range(0, len(text), size - overlap):
        chunks.append(text[i:i + size])
    return chunks

def generate_apa7_citation(filepath: str) -> str:
    try:
        mtime = os.path.getmtime(filepath)
        year = datetime.datetime.fromtimestamp(mtime).strftime('%Y')
        author = os.environ.get('USER', 'Author')
    except Exception:
        year = "n.d."
        author = "Unknown"

    filename = os.path.basename(filepath)
    mime_type, _ = mimetypes.guess_type(filepath)
    
    fmt = "Document"
    if mime_type:
        if 'pdf' in mime_type:
            fmt = "PDF document"
        elif 'spreadsheet' in mime_type or 'excel' in mime_type:
            fmt = "Excel spreadsheet"
        elif 'presentation' in mime_type or 'powerpoint' in mime_type:
            fmt = "PowerPoint presentation"
        elif 'wordprocessing' in mime_type or 'word' in mime_type:
            fmt = "Word document"
        elif 'audio' in mime_type:
            fmt = "Audio metadata"
        elif 'video' in mime_type:
            fmt = "Video metadata"
        elif 'text' in mime_type:
            fmt = "Text file"
        elif 'image' in mime_type:
            fmt = "Image metadata"

    return f"{author}. ({year}). *{filename}* [{fmt}]. Local File Index. Retrieved from file://{filepath}"

def main():
    parser = argparse.ArgumentParser(description="Deterministic Hybrid Local Search.")
    parser.add_argument("query", help="The search query.")
    parser.add_argument("--jsonIO", action="store_true", help="Output results in JSON format.")
    parser.add_argument("--ollamaExternal", type=str, default=None, help="Custom Ollama server address.")
    parser.add_argument("--ollamaHost", type=str, default=None, help="Custom Ollama host address.")
    args = parser.parse_args()

    global OLLAMA_BASE_URL, OLLAMA_EMBED_ENDPOINT
    host = args.ollamaHost or args.ollamaExternal
    if host:
        OLLAMA_BASE_URL = host if host.startswith("http") else f"http://{host}"
        OLLAMA_EMBED_ENDPOINT = f"{OLLAMA_BASE_URL}/api/embed"

    load_cache()

    if not ensure_ollama_running():
        sys.exit(1)

    if args.jsonIO:
        print(json.dumps({"phase": 1, "status": "start", "query": args.query}), flush=True)
    else:
        print(f"[*] Querying massive recoll dump for '{args.query}'...", file=sys.stderr)
    
    # Phase 1: Lexical Filter (Recoll TF-IDF -> Top 10)
    t1_start = time.perf_counter()
    top_10_files = get_file_paths_from_massive_dump(args.query, TOP_FILES_TO_PROCESS)
    t1_end = time.perf_counter()
    
    if not top_10_files:
        if args.jsonIO:
            print(json.dumps({"phase": 1, "status": "no_results"}), flush=True)
        else:
            print("❌ No files found in the index.")
        return
        
    if args.jsonIO:
        phase1_results = []
        for path in top_10_files:
            phase1_results.append({
                "path": path,
                "citation": generate_apa7_citation(path)
            })
        print(json.dumps({"phase": 1, "status": "complete", "results": phase1_results, "time_ms": (t1_end - t1_start)*1000}), flush=True)
    else:
        print(f"[*] Lexical Filter isolated Top {len(top_10_files)} documents in {(t1_end - t1_start)*1000:.2f} ms.", file=sys.stderr)

    query_emb = get_embedding(args.query)
    if query_emb is None:
        return

    all_chunks = []
    if not args.jsonIO:
        print(f"[*] Executing Python extraction for top {len(top_10_files)} files...", file=sys.stderr)
    
    for path in top_10_files:
        text = extract_content_via_python(path)[:MAX_CHARS_PER_FILE]
        if not text.strip():
            continue
        
        chunks = chunk_text(text)
        for chunk in chunks:
            all_chunks.append({"path": path, "text": chunk})

    if not args.jsonIO:
        print(f"[*] Filtering {len(all_chunks)} chunks against threshold {RANK_THRESHOLD}...", file=sys.stderr)
    
    # Phase 2: Semantic Chunking (Ollama)
    t2_start = time.perf_counter()
    
    chunk_scores = []
    seen_hashes = set()
    
    for item in all_chunks:
        h = hashlib.sha256(item['text'].encode('utf-8')).hexdigest()
        if h in seen_hashes:
            continue
        seen_hashes.add(h)
        
        c_emb = get_embedding(item['text'])
        if c_emb is not None:
            score = cosine_similarity(query_emb, c_emb)
            # Strict Relevance Cutoff
            if score >= RANK_THRESHOLD:
                chunk_scores.append((score, item))

    chunk_scores.sort(key=lambda x: x[0], reverse=True)
    final_results = chunk_scores[:TOP_CHUNKS_TO_RETURN]

    t2_end = time.perf_counter()
    phase2_ms = (t2_end - t2_start) * 1000
    
    if args.jsonIO:
        phase2_results = []
        for score, res in final_results:
            phase2_results.append({
                "score": float(score),
                "path": res['path'],
                "text": res['text'],
                "citation": generate_apa7_citation(res['path'])
            })
        print(json.dumps({"phase": 2, "status": "complete", "results": phase2_results, "time_ms": phase2_ms}), flush=True)
    else:
        print(f"   ⏱️ Chunk Embedding & Filtering completed in: {phase2_ms:.2f} ms", file=sys.stderr)

        # --- Markdown Output ---
        print("\n# Local Search Results (Threshold Filtered)", flush=True)
        print(f"*Query: {args.query} (Cutoff: {RANK_THRESHOLD})*\n", flush=True)
        print("> ℹ️ Note: If a tool suggests re-parsing a document, it may be an **Invalid trigger**. Refer to the provided content chunks. **Use these as your primary Reference.**\n", flush=True)

        if not final_results:
            print(f"⚠️ No chunks met the strict relevance threshold of {RANK_THRESHOLD}.", flush=True)
        else:
            for i, (score, res) in enumerate(final_results):
                apa_citation = generate_apa7_citation(res['path'])
                
                print(f"## {i+1}. Result (Score: {score:.4f})", flush=True)
                print(f"**Citation:** {apa_citation}\n", flush=True)
                print("### Summary Chunk (Raw 512 Chars)", flush=True)
                print("```text", flush=True)
                print(f"{res['text']}", flush=True)
                print("```\n", flush=True)
                print("---\n", flush=True)

        print(f"---\n**Performance Metrics:**\n* Lexical Pre-Filter: `{(t1_end - t1_start)*1000:.2f} ms`\n* Semantic Filtering: `{phase2_ms:.2f} ms`\n", flush=True)

    # Flush RAM to Disk if modified
    save_cache()

if __name__ == "__main__":
    print(f"[*] Invoked: {sys.executable} {' '.join(sys.argv)}", file=sys.stderr)
    main()
