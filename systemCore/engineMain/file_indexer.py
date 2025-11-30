# file_indexer.py
import os
import re
import sys
import time
import threading
import mimetypes
import datetime
import asyncio
import platform
from sqlalchemy.orm import Session
import shutil
from loguru import logger
import hashlib # <<< Import hashlib for MD5
from typing import Optional, Tuple, List, Any, Dict  # <<< ADD Optional and List HERE
import json
from cortex_backbone_provider import CortexEngine
from CortexConfiguration import * # Ensure this includes the SQLite DATABASE_URL and all prompts/models
import base64
from io import BytesIO
from PIL import Image # Requires Pillow: pip install Pillow
#from langchain_community.vectorstores import Chroma # Add Chroma import
from langchain_chroma import Chroma
import argparse # For __main__
import shutil
import subprocess
import fitz
import pytesseract
from chromadb.config import Settings

from langchain_core.messages import HumanMessage
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.output_parsers import StrOutputParser

# Optional imports - handle gracefully if libraries are missing

try:
    import pillow_avif # This plugin automatically registers AVIF support with Pillow
    AVIF_SUPPORT_AVAILABLE = True
    logger.info("✅ pillow-avif-plugin found. On-the-fly AVIF/HEIC conversion for OCR is enabled.")
except ImportError:
    AVIF_SUPPORT_AVAILABLE = False
    logger.warning("⚠️ pillow-avif-plugin not installed. On-the-fly AVIF/HEIC conversion for OCR is disabled. Run: pip install pillow-avif-plugin")

try:
    import pypdf
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False
    logger.warning("pypdf not installed. PDF text extraction disabled.")
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not installed. DOCX text extraction disabled.")
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False
    logger.warning("openpyxl not installed. XLSX text extraction disabled.")
try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not installed. PPTX text extraction disabled.")

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("pdf2image or poppler not installed. PDF -> Image conversion disabled.")


PIX2TEX_AVAILABLE = False
#Let's comment this first because it's having error or crashing the program
"""try:
    # Assuming the package installed from GitHub is importable as 'pix2tex'
    from pix2tex import cli as pix2tex
    PIX2TEX_AVAILABLE = True
    logger.info("Pix2Tex (Latex-OCR) library imported successfully.")
except ImportError:
    logger.warning("pix2tex library not found. Latex-OCR functionality will be disabled.")
    pix2tex = None # Define as None to avoid Unresolved reference errors
"""
pix2tex = None

# Database imports
DATABASE_AVAILABLE = False
try:
    # This is the normal path. All real classes are imported.
    from database import SessionLocal, FileIndex, add_interaction, init_db, IndexStatusEnum
    from sqlalchemy import update, select
    from sqlalchemy.exc import SQLAlchemyError
    DATABASE_AVAILABLE = True
    logger.trace("Successfully imported real database components.")

except ImportError:
    # This is the fallback path for the static analyzer.
    logger.critical("❌ Failed to import database components. Indexer functionality will be disabled.")

    # --- NEW: A dummy class to mimic an enum member ---
    # This class has the .value attribute the linter is looking for.
    class _DummyEnumMember:
        def __init__(self, value: str):
            self.value = value

    class SessionLocal:
        def __call__(self): return None

    # This dummy class now uses the _DummyEnumMember to satisfy the linter.
    class FileIndex:
        def __init__(self):
            self.id: int = 0
            self.file_path: str = ""
            self.file_name: str = ""
            self.size_bytes: int = 0
            self.last_modified_os: object = None
            self.last_indexed_db: object = None
            # --- MODIFIED LINE ---
            self.index_status = _DummyEnumMember("dummy_status")
            self.indexed_content: str = ""
            self.embedding_json: str = ""
            self.md5_hash: str = ""
            self.processing_error: str = ""
            self.latex_representation: str = ""
            self.latex_explanation: str = ""
            self.vlm_processing_status: str = ""

    # Dummy Enum to satisfy type hints and imports
    class IndexStatusEnum:
        pending = _DummyEnumMember("pending")
        error_permission = _DummyEnumMember("error_permission")
        # You can add other statuses here if needed for other parts of the code

    # Dummy functions and classes
    def add_interaction(*args, **kwargs): pass
    def init_db(): pass
    class SQLAlchemyError(Exception): pass
    def update(*args, **kwargs): return None
    def select(*args, **kwargs): return None

class TaskInterruptedException(Exception):
    """Custom exception raised when a lower-priority task is interrupted by a higher-priority one."""
    pass

# --- NEW: Import the custom lock ---

from shared_state import TaskInterruptedException, server_is_busy_event

from priority_lock import ELP0, ELP1 # Ensure these are imported
interruption_error_marker = "Worker task interrupted by higher priority request" # Define consistently

# --- NEW: Module-level lock and event for initialization ---
_file_index_vs_init_lock = threading.Lock()
_file_index_vs_initialized_event = threading.Event() # To
global_file_index_vectorstore: Optional[Chroma] = None

# --- Constants ---
MAX_TEXT_FILE_SIZE_MB = 1
MAX_TEXT_FILE_SIZE_BYTES = MAX_TEXT_FILE_SIZE_MB * 1024 * 1024
MAX_HASH_FILE_SIZE_MB = 500 # <<< Limit size for hashing (adjustable)
MAX_HASH_FILE_SIZE_BYTES = MAX_HASH_FILE_SIZE_MB * 1024 * 1024
SCAN_INTERVAL_HOURS = 12
SCAN_INTERVAL_SECONDS = SCAN_INTERVAL_HOURS * 60 * 60
YIELD_SLEEP_SECONDS = 0.00 # Small sleep during scanning to yield CPU (put it to 0 for maximum utilization)
MD5_CHUNK_SIZE = 65536 # Read file in chunks for hashing large files

# Add common text/code extensions if mimetypes fails
TEXT_EXTENSIONS = { # Added csv
    '.txt', '.md', '.py', '.js', '.jsx', '.ts', '.tsx', '.java', '.c', '.cpp', '.h', '.hpp',
    '.cs', '.go', '.php', '.rb', '.swift', '.kt', '.kts', '.rs', '.toml', '.yaml', '.yml',
    '.json', '.xml', '.html', '.css', '.scss', '.less', '.sh', '.bash', '.zsh', '.ps1',
    '.sql', '.log', '.m', '.cmake', '.r', '.lua', '.pl', '.pm', '.t', '.scala', '.hs',
    '.config', '.conf', '.ini', '.cfg', '.env', '.dockerfile', '.gitignore', '.gitattributes',
    '.csv', '.tsv', '.tex', '.bib'
}
DOC_EXTENSIONS = {'.pdf', '.docx', '.xlsx', '.pptx'}
OFFICE_EXTENSIONS = {'.docx', '.doc', '.xls', '.xlsx', '.pptx', '.ppt'} # Added from CortexConfiguration.py

# --- File types we want to embed ---
EMBEDDABLE_EXTENSIONS = TEXT_EXTENSIONS.union(DOC_EXTENSIONS).union({'.csv'}) # Explicitly add csv here if not in TEXT_EXTENSIONS


def log_worker(level: str, message: str):
    """Basic logging to stderr for the worker itself."""
    print(f"[{time.strftime('%Y-%m-%d %H:%M:%S')} AUDIO_WORKER(PID:{os.getpid()})|{level}] {message}", file=sys.stderr, flush=True)


class FileIndexer:
    """Scans the filesystem, extracts text/metadata, and updates the database index."""

    # --- MODIFIED: Accept embedding_model ---
    def __init__(self, stop_event: threading.Event, provider: CortexEngine, server_busy_event: threading.Event):
        self.stop_event = stop_event
        self.provider = provider
        self.embedding_model = provider.embeddings
        self.vlm_model = provider.get_model("vlm")
        self.latex_model = provider.get_model("latex")
        self.thread_name = "FileIndexerThread"
        self.server_busy_event = server_busy_event

        # --- FIX #1: ADD CRITICAL ERROR HANDLING FOR VECTOR STORE ---
        # Acquire the vector store instance during initialization. This is a blocking call.
        logger.info(f"🧵 {self.thread_name}: Acquiring vector store instance...")
        self.vectorstore = get_global_file_index_vectorstore()  # This can block for a long time
        if self.vectorstore is None:
            # If it returns None, the indexer cannot function. Log a critical error.
            # The thread will still start, but its main loop will likely fail,
            # which is better than crashing on a NoneType error.
            logger.critical(
                f"🔥🔥 {self.thread_name}: FAILED to acquire vector store instance. Indexer will not be able to store embeddings.")
        else:
            logger.success(f"✅ {self.thread_name}: Successfully acquired vector store instance.")

        # --- FIX #2: INITIALIZE THE TEXT SPLITTER ---
        # This is required for chunking documents before embedding.
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

        # Consolidated Logging for models
        emb_model_info = getattr(self.embedding_model, 'model', getattr(self.embedding_model, 'model_name', type(
            self.embedding_model).__name__)) if self.embedding_model else "Not Available"
        vlm_model_info = getattr(self.vlm_model, 'model', getattr(self.vlm_model, 'model_name', type(
            self.vlm_model).__name__)) if self.vlm_model else "Not Available"
        latex_model_info = getattr(self.latex_model, 'model', getattr(self.latex_model, 'model_name', type(
            self.latex_model).__name__)) if self.latex_model else "Not Available"

        logger.info(f"🧵 {self.thread_name} Initialized with:")
        logger.info(f"   - Embedding Model: {emb_model_info}")
        logger.info(f"   - VLM (Analysis) Model: {vlm_model_info}")
        logger.info(f"   - LaTeX (Refinement) Model: {latex_model_info}")
        logger.info(
            f"   - Text Splitter: Chunk Size={self.text_splitter._chunk_size}, Overlap={self.text_splitter._chunk_overlap}")

        if not self.embedding_model: logger.warning("⚠️ Embeddings disabled.")
        if not self.vlm_model or not self.latex_model: logger.warning(
            "⚠️ VLM->LaTeX processing disabled (one or both models missing).")

        # Initialize LatexOCR model once (this part is fine as it is)
        self.latex_ocr_model = None
        if PIX2TEX_AVAILABLE and pix2tex:
            try:
                logger.info("Initializing LatexOCR model for file indexer...")
                self.latex_ocr_model = pix2tex.LatexOCR()
                logger.success("LatexOCR model initialized successfully.")
            except Exception as e_pix2tex_init:
                logger.error(f"Failed to initialize LatexOCR model: {e_pix2tex_init}")
                self.latex_ocr_model = None

    def _wait_if_server_busy(self, check_interval=0.5, log_wait=True):
        """Checks the busy event and sleeps if set."""
        waited = False
        while self.server_busy_event.is_set():
            if not waited and log_wait: # Log only once per wait period
                logger.info(f"🚦 {self.thread_name}: Server is busy, pausing indexing...")
            waited = True
            # Check stop event frequently while waiting
            if self.stop_event.wait(timeout=check_interval):
                 logger.info(f"🛑 {self.thread_name}: Stop signaled while waiting for server.")
                 return True # Indicate stop was requested
        if waited and log_wait:
             logger.info(f"� {self.thread_name}: Server is free, resuming indexing.")
        return False # Indicate processing can continue

    def _get_latex_from_image(self, image: Image.Image) -> Optional[str]:
        """Uses pix2tex LatexOCR to extract math LaTeX from a PIL Image."""
        if not self.latex_ocr_model:
            logger.warning("LatexOCR model not available, skipping math extraction.")
            return None

        try:
            log_worker("TRACE", "Extracting math using LatexOCR model...")
            # The model call is synchronous/blocking
            math_latex = self.latex_ocr_model(image)
            if math_latex and math_latex.strip():
                log_worker("DEBUG", f"LatexOCR extracted math: {math_latex.strip()[:100]}...")
                return math_latex.strip()
            return None  # Return None if OCR produces no text
        except Exception as e_latex_ocr:
            log_worker("ERROR", f"Error during LatexOCR processing: {e_latex_ocr}")
            return f"[LatexOCR Error: {e_latex_ocr}]"

    def _extract_text_with_ocr_fallback(self, file_path: str, file_ext: str) -> str:
        """
        Extracts text from a file. For PDFs, it checks for a text layer and
        falls back to OCR. For standard image types, it performs OCR directly,
        with a fallback to convert unsupported formats like AVIF/HEIC on the fly.
        """
        log_prefix = f"TextExtract|{os.path.basename(file_path)[:15]}"

        try:
            # --- PDF Processing (Unchanged) ---
            if file_ext == '.pdf':
                doc: fitz.Document = fitz.open(file_path)
                full_text = ""
                is_image_based = True
                page: fitz.Page
                for page in doc:
                    text_from_page = page.get_text("text")
                    if text_from_page and len(text_from_page.strip()) > 20:
                        is_image_based = False
                        break
                if not is_image_based:
                    logger.trace(f"{log_prefix}: PDF has text layer, extracting directly.")
                    full_text = "\n".join([page.get_text() for page in doc])
                else:
                    logger.info(f"{log_prefix}: PDF has no text layer, performing OCR...")
                    ocr_texts = []
                    for page_num, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=300)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        try:
                            ocr_texts.append(pytesseract.image_to_string(img))
                        except pytesseract.TesseractNotFoundError:
                            logger.error(f"{log_prefix}: Tesseract OCR engine not found.")
                            return "[OCR Error: Tesseract engine not found]"
                    full_text = "\n\n--- Page Break ---\n\n".join(ocr_texts)
                doc.close()
                return full_text

            # --- Image Processing (NEW LOGIC) ---
            elif file_ext in OCR_TARGET_EXTENSIONS:
                logger.info(f"{log_prefix}: Performing OCR on image file: {file_path}")
                try:
                    # Attempt direct OCR first. This works for JPG, PNG, etc.
                    return pytesseract.image_to_string(Image.open(file_path))
                except pytesseract.TesseractNotFoundError:
                    logger.error(f"{log_prefix}: Tesseract OCR engine not found.")
                    return "[OCR Error: Tesseract engine not found]"
                except Exception as e:
                    # Check if this is a format error that we can handle with conversion.
                    if "Unsupported image format" in str(e) or "cannot identify image file" in str(e):
                        logger.warning(
                            f"{log_prefix}: Unsupported format '{file_ext}'. Attempting on-the-fly conversion...")

                        if not AVIF_SUPPORT_AVAILABLE:
                            err_msg = f"Cannot convert '{file_path}'. The required AVIF/HEIC plugin ('pillow-avif-plugin') is not installed."
                            logger.error(f"{log_prefix}: {err_msg}")
                            return f"[OCR Error: {err_msg}]"

                        # If the plugin is available, try the conversion logic.
                        try:
                            # Load the unsupported image (this should work now thanks to the plugin)
                            unsupported_image = Image.open(file_path)

                            # Convert to PNG in memory for maximum compatibility.
                            # The .convert("RGB") is crucial for formats that might have alpha channels (like PNG/AVIF).
                            png_buffer = BytesIO()
                            unsupported_image.convert("RGB").save(png_buffer, format="PNG")
                            png_buffer.seek(0)  # Rewind the buffer to the beginning

                            # Now, run OCR on the in-memory PNG data
                            logger.info(f"{log_prefix}: Conversion successful. Performing OCR on in-memory PNG.")
                            return pytesseract.image_to_string(Image.open(png_buffer))
                        except Exception as conversion_error:
                            err_msg = f"On-the-fly conversion failed for '{file_path}': {conversion_error}"
                            logger.error(f"{log_prefix}: {err_msg}")
                            return f"[OCR Error: {err_msg}]"
                    else:
                        # Re-raise any other unexpected errors.
                        raise e

        except Exception as e:
            logger.error(f"{log_prefix}: Failed to extract text/OCR from '{file_path}': {e}")
            return f"[Error extracting text/OCR: {e}]"

        return ""  # Fallback for unhandled file types

    def _convert_office_to_pdf(self, office_path: str) -> Optional[str]:
        """
        Converts an Office document (docx, xlsx, pptx) to a temporary PDF file
        using unoconv (requires LibreOffice installed).
        Returns the path to the temporary PDF file or None on failure.
        """
        if not shutil.which("unoconv"):
            logger.error(f"unoconv command not found. Cannot convert Office file: {office_path}")
            return None

        temp_pdf_dir = os.path.join(os.path.dirname(office_path), "temp_office_conversions")
        os.makedirs(temp_pdf_dir, exist_ok=True)
        base_name = os.path.splitext(os.path.basename(office_path))[0]
        temp_pdf_path = os.path.join(temp_pdf_dir, f"{base_name}.pdf")

        logger.debug(f"Attempting to convert '{office_path}' to PDF at '{temp_pdf_path}' using unoconv...")
        try:
            # Command: unoconv -f pdf -o /path/to/output_dir /path/to/input_file
            # We specify output file directly with -o for some unoconv versions.
            # If -o expects a dir, then use:
            # process = subprocess.run(["unoconv", "-f", "pdf", "-o", temp_pdf_dir, office_path],
            process = subprocess.run(["unoconv", "-f", "pdf", "-o", temp_pdf_path, office_path],
                                     capture_output=True, text=True, timeout=120, check=False)
            if process.returncode == 0 and os.path.exists(temp_pdf_path):
                logger.info(f"Successfully converted '{office_path}' to '{temp_pdf_path}'.")
                return temp_pdf_path
            else:
                logger.error(f"unoconv conversion failed for '{office_path}'. RC: {process.returncode}")
                logger.error(f"unoconv stdout: {process.stdout.strip()}")
                logger.error(f"unoconv stderr: {process.stderr.strip()}")
                if os.path.exists(temp_pdf_path): # Clean up if partial file created
                    try: os.remove(temp_pdf_path)
                    except: pass
                return None
        except subprocess.TimeoutExpired:
            logger.error(f"unoconv conversion timed out for '{office_path}'.")
            if os.path.exists(temp_pdf_path):
                try: os.remove(temp_pdf_path)
                except: pass
            return None
        except Exception as e:
            logger.error(f"Error during unoconv conversion of '{office_path}': {e}")
            if os.path.exists(temp_pdf_path):
                try: os.remove(temp_pdf_path)
                except: pass
            return None


    def _convert_pdf_to_images(self, pdf_path: str) -> Optional[List[Any]]:
         """Converts PDF pages to PIL Image objects."""
         if not PDF2IMAGE_AVAILABLE: return None
         logger.trace(f"Converting PDF to images: {pdf_path}")
         try:
             # Use a high enough DPI for potentially small formulas, but balance size
             images = convert_from_path(pdf_path, dpi=200)
             logger.trace(f"  Converted {len(images)} pages.")
             return images
         except Exception as e:
             logger.error(f"Failed PDF to image conversion for {pdf_path}: {e}")
             return None

    def _get_initial_vlm_description(self, image: Image.Image) -> Tuple[Optional[str], Optional[str]]:
        """Calls the 'vlm' model for initial analysis/description with ELP0 priority."""
        # Check server busy/stop event first
        if self._wait_if_server_busy(): return None, "[VLM Skipped - Server Busy]"
        if self.stop_event.is_set(): return None, "[VLM Skipped - Stop Requested]"
        if not self.vlm_model: return None, "[VLM Model Unavailable]"

        logger.trace("Sending image page to VLM Model for initial analysis (Priority: ELP0)...")
        raw_description = None # Initialize
        try:
            # Prepare image data
            buffered = BytesIO(); image.save(buffered, format="PNG"); img_str = base64.b64encode(buffered.getvalue()).decode('utf-8'); image_uri = f"data:image/png;base64,{img_str}"
            image_content_part = {"type": "image_url", "image_url": {"url": image_uri}}

            # Prepare messages and chain
            messages = [HumanMessage(content=[image_content_part, {"type": "text", "text": PROMPT_VLM_INITIAL_ANALYSIS}])]
            chain = self.vlm_model | StrOutputParser()

            # Invoke the chain with ELP0 priority
            # Assuming chain.invoke supports the 'config' dictionary for priority
            raw_description = chain.invoke(messages, config={'priority': ELP0})

            logger.trace(f"  VLM initial analysis raw response length: {len(raw_description)}")
            return raw_description, None # Return description, no error string if successful

        except Exception as e:
            # --- Interruption Handling ---
            if interruption_error_marker in str(e):
                logger.warning(f"🚦 VLM initial analysis INTERRUPTED by ELP1.")
                return None, "[VLM Interrupted]" # Specific marker for interruption
            # --- End Interruption Handling ---
            else:
                # Handle other errors
                logger.error(f"VLM initial analysis call failed: {e}", exc_info=True)
                return None, f"[VLM Analysis Error: {e}]" # Generic error marker

    # --- MODIFIED: _refine_to_latex_tikz (Adds ELP0 & Interruption Handling) ---

    # --- NEW HELPER METHOD: _refine_to_latex_tikz ---
    def _refine_to_latex_tikz(self, image: Image.Image, initial_analysis: str) -> Tuple[Optional[str], Optional[str]]:
        """Calls the 'latex' model to refine VLM analysis into LaTeX/TikZ with ELP0 priority."""
        # Check server busy/stop event first
        if self._wait_if_server_busy(): return None, "[LaTeX Model Skipped - Server Busy]"
        if self.stop_event.is_set(): return None, "[LaTeX Model Skipped - Stop Requested]"
        if not self.latex_model: return None, "[LaTeX Model Unavailable]"

        logger.trace("Sending image and initial analysis to LaTeX Model for refinement (Priority: ELP0)...")
        response_markdown = None # Initialize
        try:
            # Prepare image data
            buffered = BytesIO(); image.save(buffered, format="PNG"); img_str = base64.b64encode(buffered.getvalue()).decode('utf-8'); image_uri = f"data:image/png;base64,{img_str}"
            image_content_part = {"type": "image_url", "image_url": {"url": image_uri}}

            # Prepare messages and chain
            prompt_text = PROMPT_LATEX_REFINEMENT.format(initial_analysis=initial_analysis)
            messages = [HumanMessage(content=[image_content_part, {"type": "text", "text": prompt_text}])]
            chain = self.latex_model | StrOutputParser()

            # Invoke the chain with ELP0 priority
            response_markdown = chain.invoke(messages, config={'priority': ELP0})

            # --- Parsing logic (remains the same) ---
            latex_code = None; tikz_code = None; explanation = response_markdown.strip()
            try:
                latex_match = re.search(r"```(?:latex)?\s*(.*?)\s*```", response_markdown, re.DOTALL | re.IGNORECASE)
                tikz_match = re.search(r"```(?:tikz)?\s*(.*?)\s*```", response_markdown, re.DOTALL | re.IGNORECASE)
                cleaned_response = response_markdown
                if latex_match: latex_code = latex_match.group(1).strip(); cleaned_response = cleaned_response.replace(latex_match.group(0), "", 1); logger.trace("  LaTeX Model extracted LaTeX.")
                if tikz_match: tikz_code = tikz_match.group(1).strip(); cleaned_response = cleaned_response.replace(tikz_match.group(0), "", 1); logger.trace("  LaTeX Model extracted TikZ.")
                explanation = cleaned_response.strip()
                if not explanation and (latex_code or tikz_code): explanation = "(Code extracted, no separate explanation provided by model)"
            except Exception as parse_e: logger.warning(f"Could not parse refined LaTeX/TikZ response: {parse_e}")
            # Combine LaTeX and TikZ for storage
            final_latex_repr = latex_code
            if tikz_code: final_latex_repr = (final_latex_repr or "") + "\n\n% --- TikZ Code ---\n" + tikz_code
            # --- End Parsing Logic ---

            return final_latex_repr, explanation # Return parsed results

        except Exception as e:
            # --- Interruption Handling ---
            if interruption_error_marker in str(e):
                logger.warning(f"🚦 LaTeX refinement INTERRUPTED by ELP1.")
                return None, "[LaTeX Refinement Interrupted]" # Specific marker
            # --- End Interruption Handling ---
            else:
                # Handle other errors
                logger.error(f"LaTeX model refinement call failed: {e}", exc_info=True)
                return None, f"[LaTeX Refinement Error: {e}]" # Generic error marker

    # --- NEW: MD5 Hashing Helper ---
    def _calculate_md5(self, file_path: str, size_bytes: Optional[int]) -> Optional[str]:
        """Calculates the MD5 hash of a file, reading in chunks."""
        if size_bytes is None or size_bytes < 0:
             logger.warning(f"Cannot calculate MD5 for {file_path}: Unknown size.")
             return None
        if size_bytes > MAX_HASH_FILE_SIZE_BYTES:
             logger.trace(f"Skipping MD5 hash (file too large: {size_bytes / 1024 / 1024:.1f} MB > {MAX_HASH_FILE_SIZE_MB} MB): {file_path}")
             return None # Skip hashing very large files for performance

        logger.trace(f"Calculating MD5 hash for: {file_path}")
        hash_md5 = hashlib.md5()
        try:
            with open(file_path, "rb") as f:
                while True:
                    if self.stop_event.is_set():
                         logger.info(f"MD5 calculation interrupted by stop signal: {file_path}")
                         return None # Stop hashing
                    chunk = f.read(MD5_CHUNK_SIZE)
                    if not chunk:
                        break
                    hash_md5.update(chunk)
            hex_digest = hash_md5.hexdigest()
            logger.trace(f"MD5 calculated for {file_path}: {hex_digest}")
            return hex_digest
        except FileNotFoundError:
            logger.warning(f"MD5 failed: File not found at {file_path}")
            return None
        except PermissionError:
            raise # Re-raise permission errors to be handled by _process_file
        except Exception as e:
            logger.error(f"MD5 calculation failed for {file_path}: {e}")
            return None # Indicate hashing failure
    # --- End MD5 Helper ---

    def _get_root_paths(self) -> list[str]:
        """Returns a list of root directories to scan based on the OS."""
        paths = []
        if sys.platform.startswith("win"):
            import string
            logger.info("Detected Windows platform.")
            for drive in string.ascii_uppercase:
                path = f"{drive}:\\"
                if os.path.exists(path):
                    paths.append(path)
                    logger.debug(f"Adding drive for scanning: {path}")
        elif sys.platform.startswith("darwin"):
            logger.info("Detected macOS platform.")
            paths = ["/"] # Start from root on macOS
            # Add /Volumes explicitly if needed, but root should cover it
            # if os.path.exists("/Volumes"): paths.append("/Volumes")
        elif sys.platform.startswith("linux"):
            logger.info("Detected Linux platform.")
            paths = ["/"] # Start from root on Linux
            # Common mount points often under root or /mnt
            # if os.path.exists("/mnt"): paths.append("/mnt")
            # if os.path.exists("/media"): paths.append("/media")
        else:
            logger.warning(f"Unsupported platform '{sys.platform}'. Using '/' as default root.")
            paths = ["/"]

        logger.info(f"Root paths for scanning: {paths}")
        return paths

    def _get_file_metadata(self, file_path: str) -> Dict[str, Any]:
        """Gets file metadata and returns it as a dictionary."""
        try:
            stat_result = os.stat(file_path)
            return {
                'size_bytes': stat_result.st_size,
                'last_modified_os': datetime.datetime.fromtimestamp(stat_result.st_mtime, tz=datetime.timezone.utc),
                'mime_type': mimetypes.guess_type(file_path)[0]
            }
        except (FileNotFoundError, PermissionError) as e:
            logger.warning(f"Metadata read error for {file_path}: {e}")
            return {'size_bytes': None, 'last_modified_os': None, 'mime_type': None}

    def _extract_text(self, file_path: str, size_bytes: int) -> Optional[str]:
        """Attempts to read content from presumed text files."""
        if size_bytes > MAX_TEXT_FILE_SIZE_BYTES:
            logger.trace(f"Skipping text read (too large: {size_bytes / 1024 / 1024:.1f} MB): {file_path}")
            return None # Indicate skipped due to size

        logger.trace(f"Attempting text read: {file_path}")
        try:
            # Try UTF-8 first, ignore errors for robustness
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except IsADirectoryError:
             logger.warning(f"Attempted to read directory as text file: {file_path}")
             return None # Should not happen if called correctly, but handle anyway
        except PermissionError:
            raise # Re-raise permission errors to be handled by _process_file
        except Exception as e:
            logger.warning(f"Text read failed for {file_path}: {e}")
            # Optionally try other encodings here if needed
            return None # Indicate read failure

    def _extract_pdf(self, file_path: str) -> Optional[str]:
        """Extracts text from a PDF file using pypdf."""
        if not PYPDF_AVAILABLE: return None
        logger.trace(f"Attempting PDF extraction: {file_path}")
        text_content = ""
        try:
            reader = pypdf.PdfReader(file_path)
            num_pages = len(reader.pages)
            logger.trace(f"  PDF has {num_pages} pages.")
            for i, page in enumerate(reader.pages):
                if self.stop_event.is_set(): return "[Indexing stopped]"
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_content += page_text + "\n"
                except Exception as page_err:
                    # Log error for specific page but continue
                    logger.warning(f"Error extracting text from page {i+1} of {file_path}: {page_err}")
            return text_content.strip() if text_content else None
        except pypdf.errors.PdfReadError as pe:
             logger.warning(f"Failed to read PDF (possibly encrypted or corrupted): {file_path} - {pe}")
             return None
        except PermissionError:
             raise # Re-raise permission errors
        except Exception as e:
            logger.error(f"PDF extraction failed for {file_path}: {e}")
            return None

    def _extract_docx(self, file_path: str) -> Optional[str]:
        """Extracts text from a DOCX file using python-docx."""
        if not DOCX_AVAILABLE: return None
        logger.trace(f"Attempting DOCX extraction: {file_path}")
        try:
            document = docx.Document(file_path)
            full_text = [para.text for para in document.paragraphs]
            return '\n'.join(full_text).strip() if full_text else None
        except PermissionError:
            raise # Re-raise permission errors
        except Exception as e:
            logger.error(f"DOCX extraction failed for {file_path}: {e}")
            return None

    def _extract_xlsx(self, file_path: str) -> Optional[str]:
        """Extracts text from an XLSX file using openpyxl."""
        if not OPENPYXL_AVAILABLE: return None
        logger.trace(f"Attempting XLSX extraction: {file_path}")
        text_content = ""
        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True) # data_only for values
            for sheetname in workbook.sheetnames:
                if self.stop_event.is_set(): return "[Indexing stopped]"
                sheet = workbook[sheetname]
                text_content += f"--- Sheet: {sheetname} ---\n"
                for row in sheet.iter_rows():
                    row_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                    text_content += "\t".join(row_values) + "\n"
                text_content += "\n"
            return text_content.strip() if text_content else None
        except PermissionError:
            raise # Re-raise permission errors
        except Exception as e:
            logger.error(f"XLSX extraction failed for {file_path}: {e}")
            return None

    def _extract_pptx(self, file_path: str) -> Optional[str]:
        """Extracts text from a PPTX file using python-pptx."""
        if not PPTX_AVAILABLE: return None
        logger.trace(f"Attempting PPTX extraction: {file_path}")
        text_content = ""
        try:
            presentation = pptx.Presentation(file_path)
            for i, slide in enumerate(presentation.slides):
                if self.stop_event.is_set(): return "[Indexing stopped]"
                text_content += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            return text_content.strip() if text_content else None
        except PermissionError:
            raise # Re-raise permission errors
        except Exception as e:
            logger.error(f"PPTX extraction failed for {file_path}: {e}")
            return None

    def _scan_directory(self, root_path: str, db_session: Session, loop: asyncio.AbstractEventLoop):
        """
        Phase 1: Walks through a directory and processes files using _process_file_phase1,
        skipping OS/system dirs/files and hidden dot files/dirs.
        Reports progress periodically.
        """
        logger.info(f"🔬 Starting Phase 1 Scan Cycle for root: {root_path}")
        total_processed_this_root_scan = 0
        total_errors_this_root_scan = 0
        last_report_time = time.monotonic()
        files_since_last_report = 0
        errors_since_last_report = 0
        report_interval_seconds = 60  # Log progress every minute

        # --- Comprehensive Directory Exclusion Configuration ---
        # globals().get("ROOT_DIR", ".") attempts to get ROOT_DIR if it's a global in file_indexer.py
        # If not, it defaults to ".", which means paths like "./build" would be relative to CWD
        # This is okay if file_indexer.py is run from the project root.
        # For more robustness, ROOT_DIR could be passed to FileIndexer or determined from __file__.
        # For now, assuming ROOT_DIR is accessible or paths are intended to be relative to CWD if it's not.
        current_script_dir_for_project_paths = os.path.dirname(os.path.dirname(
            os.path.abspath(__file__)))  # Assuming file_indexer is in engineMain, and ROOT_DIR is one level up

        common_system_dirs_raw = {
            # macOS Specific
            "/System", "/Library", os.path.expanduser("~/Library"), "/private", "/cores", "/Network",
            "/Applications", "/.Spotlight-V100", "/.fseventsd", "/.Trashes", "/usr", "/opt", "/System", "/bin", "/cores", "/dev", "/etc", "/private", "/tmp", "/var", os.path.expanduser("~/.Trash"),
            "/Volumes",  # Often external drives or system volumes, can be very large.

            # Linux Specific
            "/proc", "/sys", "/dev", "/run", "/lost+found", "/mnt", "/media", "/usr", "/boot", "/srv",
            "/snap", "/var/lib/snapd", "/var/lib/flatpak",
            os.path.expanduser("~/.cache"), os.path.expanduser("~/.config"), os.path.expanduser("~/.local/share"),
            # Standard Linux system dirs often symlinked or part of /usr, /etc, /var (already listed)

            # Windows Specific (using environment variables)
            "%SystemRoot%", "%WinDir%", "%ProgramFiles%", "%ProgramFiles(x86)%",
            "%ProgramData%", "%LocalAppData%", "%AppData%",
            "%UserProfile%\\AppData\\Local\\Temp", "%SystemDrive%\\$Recycle.Bin",
            "%SystemDrive%\\System Volume Information", "%TEMP%", "%TMP%",

            # Common Development & Build Folders (Relative Names - will be lowercased)
            "node_modules", "__pycache__", ".venv", "venv", "env", ".env",
            "target", "build", "dist", "out", "bin", "obj",  # Note: "bin" and "lib" are very generic
            ".git", ".hg", ".svn", ".gradle", ".idea", ".vscode", ".project",
            ".pytest_cache", ".mypy_cache", ".ruff_cache", "site-packages", "dist-packages",

            # Common Cache & Temporary Directory Names (Relative Names)
            "cache", "caches", "tmp", "temp", "logs", "log",  # "logs" is also your project log dir

            # Project Specific (ensure ROOT_DIR is correctly resolved for these)
            "staticmodelpool", "llama-cpp-python_build", "stable-diffusion-cpp-python_build",
            "pywhispercpp_build", "systemCore", "engineMain", "backend-service",
            "frontend-face-zephyrine", "zephyrineCondaVenv", "db_snapshots",
            "temp_audio_worker_files",
            # Explicitly add project's own log directory if it's not already covered
            os.path.join(current_script_dir_for_project_paths, "logs"),
            # Add project's build directory if it has one at the root
            os.path.join(current_script_dir_for_project_paths, "build"),
        }


        #Override to be null nothing escapes from the eyes and everything must be read for the knowledge (Only for debugging purposes)
        #common_system_dirs_raw = {}

        absolute_skip_dirs_resolved = set()
        relative_skip_dir_names_lower = set()

        for p_raw_item in common_system_dirs_raw:
            p_expanded_item = p_raw_item
            IS_WINDOWS = os.name == 'nt'
            if "%" in p_raw_item and (IS_WINDOWS or platform.system() == "Windows"):
                p_expanded_item = os.path.expandvars(p_raw_item)
                if p_expanded_item == p_raw_item and p_raw_item.count('%') >= 2:  # Unexpanded var
                    logger.trace(f"Path with '%' not expanded: '{p_raw_item}'. Treating as relative name.")

            if p_expanded_item.startswith('~'):
                p_expanded_item = os.path.expanduser(p_expanded_item)

            if os.path.isabs(p_expanded_item):
                absolute_skip_dirs_resolved.add(os.path.normcase(os.path.realpath(p_expanded_item)))
            else:
                relative_skip_dir_names_lower.add(p_expanded_item.lower())

        logger.debug(
            f"Resolved {len(absolute_skip_dirs_resolved)} absolute skip paths. Examples: {list(absolute_skip_dirs_resolved)[:3]}")
        logger.debug(
            f"Resolved {len(relative_skip_dir_names_lower)} relative skip names. Examples: {list(relative_skip_dir_names_lower)[:3]}")

        files_to_skip_lower = {
            '.ds_store', 'thumbs.db', 'desktop.ini', '.localized',
            '.bash_history', '.zsh_history', 'ntuser.dat', '.swp', '.swo',
            'pagefile.sys', 'hiberfil.sys', '.volumeicon.icns', '.cfusertextencoding',
            '.traceroute.log', '.bash_sessions_disable'
        }
        # --- End Skip Logic Init ---

        try:
            for current_dir, dirnames, filenames in os.walk(root_path, topdown=True, onerror=None):
                #logger.debug(f"SCANNER_WALK_DEBUG: Entering directory: {current_dir}")
                if self.stop_event.is_set():
                    logger.info(f"Phase 1 Scan interrupted by stop signal in {current_dir}")
                    break

                norm_current_dir_for_abs_match = os.path.normcase(os.path.realpath(current_dir))
                current_dir_basename_lower_for_rel_match = os.path.basename(current_dir).lower()
                should_skip_dir = False

                # 1. Check against resolved absolute paths and their subdirectories
                for skip_abs_path in absolute_skip_dirs_resolved:
                    if norm_current_dir_for_abs_match == skip_abs_path or \
                            norm_current_dir_for_abs_match.startswith(skip_abs_path + os.sep):
                        should_skip_dir = True
                        break

                if not should_skip_dir and current_dir_basename_lower_for_rel_match in relative_skip_dir_names_lower:
                    should_skip_dir = True

                if not should_skip_dir and os.path.basename(current_dir).startswith(
                        '.'):  # Check original name for hidden
                    should_skip_dir = True

                if should_skip_dir:
                    logger.trace(f"Phase 1 Skipping directory: {current_dir}")
                    dirnames[:] = []
                    filenames[:] = []
                    continue

                dirnames[:] = [d for d in dirnames if not d.startswith('.')]  # Prune hidden subdirs from recursion

                current_dir_file_errors = 0
                for filename in filenames:
                    if self.stop_event.is_set(): break
                    if filename.startswith('.'): continue
                    if filename.lower() in files_to_skip_lower: continue

                    file_path = os.path.join(current_dir, filename)
                    file_processed_this_iter = False  # Renamed for clarity
                    file_errored_this_iter = False  # Renamed for clarity

                    try:
                        if os.path.islink(file_path): continue
                        if not os.path.isfile(file_path): continue

                        loop.run_until_complete(
                            self._process_file_phase1(file_path, db_session)
                        )
                        file_processed_this_iter = True
                    except PermissionError:
                        logger.warning(f"Phase 1 Permission denied processing file: {file_path}")
                        file_errored_this_iter = True
                        try:
                            existing_rec = db_session.query(FileIndex).filter(FileIndex.file_path == file_path).first()
                            err_vals = {'index_status': 'error_permission',
                                        'processing_error': "Permission denied during scan.",
                                        'last_indexed_db': datetime.datetime.now(datetime.timezone.utc)}
                            if existing_rec:
                                if existing_rec.index_status != 'error_permission': db_session.execute(
                                    update(FileIndex).where(FileIndex.id == existing_rec.id).values(**err_vals))
                            else:
                                db_session.add(
                                    FileIndex(file_path=file_path, file_name=filename, **err_vals))  
                            db_session.commit()
                        except Exception as db_perm_err:
                            logger.error(
                                f"Failed to log perm error for {file_path}: {db_perm_err}"); db_session.rollback()
                    except Exception as walk_process_err:
                        logger.error(
                            f"Phase 1 Error during _process_file_phase1 call for {file_path}: {walk_process_err}",
                            exc_info=True)
                        file_errored_this_iter = True
                    finally:
                        if file_processed_this_iter: total_processed_this_root_scan += 1; files_since_last_report += 1
                        if file_errored_this_iter: total_errors_this_root_scan += 1; errors_since_last_report += 1; current_dir_file_errors += 1

                        current_time_monotonic = time.monotonic()
                        if current_time_monotonic - last_report_time >= report_interval_seconds:
                            rate = files_since_last_report / report_interval_seconds if report_interval_seconds > 0 else files_since_last_report
                            logger.info(
                                f"⏳ [Phase 1 Report] In '{os.path.basename(root_path)}' last {report_interval_seconds}s: {files_since_last_report} files (~{rate:.1f}/s), {errors_since_last_report} errors. Root Total: {total_processed_this_root_scan}")
                            last_report_time = current_time_monotonic
                            files_since_last_report = 0
                            errors_since_last_report = 0

                        if YIELD_SLEEP_SECONDS > 0: time.sleep(YIELD_SLEEP_SECONDS)

                if self.stop_event.is_set(): break
                if current_dir_file_errors > 0: logger.warning(
                    f"Encountered {current_dir_file_errors} errors processing files within: {current_dir}")

        except Exception as outer_walk_err:
            logger.error(f"Outer error during Phase 1 os.walk for {root_path}: {outer_walk_err}", exc_info=True)
            total_errors_this_root_scan += 1

        if not self.stop_event.is_set():
            logger.info(
                f"✅ Finished Phase 1 Scan for {root_path}. Total Processed this root: {total_processed_this_root_scan}, Total Errors this root: {total_errors_this_root_scan}")
        else:
            logger.warning(
                f"⏹️ Phase 1 Scan for {root_path} interrupted. Processed in this root-cycle: {total_processed_this_root_scan}, Errors: {total_errors_this_root_scan}")

    async def _embed_and_store(self, record: FileIndex, text_to_embed: str, db_session: Session, is_final_embedding: bool):
        """
        Handles the core logic of chunking, embedding, and storing vectors in ChromaDB.
        This is now called by both Phase 1 (for initial text) and Phase 2 (for enriched text).
        """
        log_prefix = f"Embed|ID:{record.id}"
        
        if not self.embedding_model:
            logger.error(f"{log_prefix}: Embedding model not available. Cannot proceed.")
            record.index_status = IndexStatusEnum.error_embedding
            record.processing_error = "Embedding model not configured."
            db_session.commit()
            return

        if not text_to_embed:
            logger.warning(f"{log_prefix}: No text content provided to embed. Marking as indexed_meta.")
            record.index_status = IndexStatusEnum.indexed_meta
            db_session.commit()
            return
            
        try:
            # Chunk the text
            chunks = self.text_splitter.split_text(text_to_embed)
            if not chunks: raise ValueError("Text splitter produced no chunks.")
            
            # Embed the chunks (as a background task)
            chunk_embeddings = await asyncio.to_thread(
                self.embedding_model.embed_documents, chunks, priority=ELP0
            )
            if not chunk_embeddings or len(chunk_embeddings) != len(chunks):
                raise ValueError("Embedding model returned empty or mismatched vectors.")

            # Update ChromaDB
            global_file_vs = get_global_file_index_vectorstore()
            if not global_file_vs:
                raise RuntimeError("Global file vector store is not available for updating.")
            
            # Always delete old entries for this file_id to ensure atomicity
            await asyncio.to_thread(global_file_vs.delete, where={"file_id": record.id})

            chunk_ids = [f"file_{record.id}_chunk_{i}" for i in range(len(chunks))]
            chunk_metadatas = [{
                "file_id": record.id,
                "file_path": record.file_path, 
                "file_name": record.file_name, 
                "chunk_index": i,
                "is_final_vlm": is_final_embedding # Add flag for search-time filtering if needed
            } for i in range(len(chunks))]

            # await asyncio.to_thread(
            #     global_file_vs.add_embeddings,
            #     embeddings=chunk_embeddings,
            #     metadatas=chunk_metadatas,
            #     ids=chunk_ids
            # )

            # --- NEW, CORRECTED CODE ---
            await asyncio.to_thread(
                global_file_vs._collection.add,  # Access the underlying collection
                embeddings=chunk_embeddings,
                metadatas=chunk_metadatas,
                ids=chunk_ids,
                documents=chunks  # It's good practice to add the document text as well
            )
            
            # Update the record's final status in the database
            record.index_status = IndexStatusEnum.indexed_complete if is_final_embedding else IndexStatusEnum.indexed_text
            record.processing_error = None # Clear previous errors
            logger.success(f"✅ {log_prefix}: Successfully embedded {len(chunks)} chunks. Final Status: {record.index_status.value}")

        except TaskInterruptedException as tie:
            logger.warning(f"🚦 {log_prefix}: Embedding INTERRUPTED: {tie}. Status will remain pending.")
            # Keep status as 'pending_embedding' to allow retry
            record.index_status = IndexStatusEnum.pending_embedding
        except Exception as e:
            logger.error(f"❌ {log_prefix}: Failed to embed or update vector store: {e}", exc_info=True)
            record.index_status = IndexStatusEnum.error_embedding
            record.processing_error = f"Embedding/Chroma update failed: {str(e)[:255]}"
        
        finally:
            db_session.commit()

    async def _embed_and_update_vector_store(self, record: FileIndex, db_session: Session):
        """
        Takes a FileIndex record, combines its text fields, chunks, embeds,
        and updates the Chroma vector store. This is the final step for any file.
        """

        log_prefix = f"Embed|{os.path.basename(record.file_path)[:15]}" #type: ignore
        logger.info(f"--> {log_prefix}: Starting final embedding process for File ID {record.id}") #type: ignore

        if not self.embedding_model:
            logger.error(
                f"{log_prefix}: Embedding model not available. Cannot complete indexing for File ID {record.id}.")
            record.index_status = 'error_embedding'
            record.processing_error = "Embedding model not available at final stage."
            db_session.commit()
            return

        # 1. Combine all available text fields into a master document
        # We add headers to give the embedding model context about the source of the text.
        master_document_text = ""
        if record.indexed_content:
            master_document_text += f"--- Raw Text Content ---\n{record.indexed_content}\n\n"
        if record.latex_explanation:  # This holds the VLM description
            master_document_text += f"--- Visual Description (VLM Analysis) ---\n{record.latex_explanation}\n\n"
        if record.latex_representation:  # This holds the LaTeX-OCR math
            master_document_text += f"--- Extracted Mathematical Notation (LaTeX-OCR) ---\n{record.latex_representation}\n\n"

        master_document_text = master_document_text.strip()

        if not master_document_text:
            logger.warning(
                f"{log_prefix}: No text content found to embed for File ID {record.id}. Marking as indexed_meta.")
            record.index_status = 'indexed_meta'
            db_session.commit()
            return

        # 2. Chunk the combined text
        try:
            # RecursiveCharacterTextSplitter is initialized in self.vector_store_init_and_load
            # Assuming self.text_splitter is available
            chunks = self.text_splitter.split_text(master_document_text)  
            logger.info(f"{log_prefix}: Split combined document into {len(chunks)} chunks.")
            if not chunks:
                raise ValueError("Text splitter produced no chunks.")
        except Exception as e_split:
            logger.error(f"{log_prefix}: Failed to split text for File ID {record.id}: {e_split}")
            record.index_status = 'error_embedding'
            record.processing_error = f"Text splitting failed: {e_split}"
            db_session.commit()
            return

        # 3. Embed the chunks
        try:
            vectorstore = get_global_file_index_vectorstore()
            # Use a lower priority for this background embedding task
            chunk_embeddings = await asyncio.to_thread(
                self.embedding_model.embed_documents, chunks, priority=ELP0  
            )
            if not chunk_embeddings or len(chunk_embeddings) != len(chunks):
                raise ValueError("Embedding model returned empty or mismatched number of vectors.")

            # 4. Update ChromaDB: Delete old entries and add new chunks
            logger.info(f"{log_prefix}: Embedding successful. Updating vector store for {len(chunks)} chunks...")

            # Delete any existing vectors for this file ID to avoid duplicates
            await asyncio.to_thread(
                self.vectorstore.delete,  
                where={"file_id": record.id}
            )

            # Create new IDs for each chunk, but use the same metadata
            chunk_ids = [f"file_{record.id}_chunk_{i}" for i in range(len(chunks))]
            chunk_metadatas = [{
                "file_id": record.id,
                "file_path": record.file_path, 
                "file_name": record.file_name, 
                "chunk_index": i
            } for i in range(len(chunks))]

            # await asyncio.to_thread(
            #     self.vectorstore.add_embeddings,  
            #     embeddings=chunk_embeddings,
            #     metadatas=chunk_metadatas,
            #     ids=chunk_ids
            # )

            # --- NEW, CORRECTED CODE ---
            await asyncio.to_thread(
                self.vectorstore._collection.add,   # Access the underlying collection
                embeddings=chunk_embeddings,
                metadatas=chunk_metadatas,
                ids=chunk_ids,
                documents=chunks  # Also add the document text
            )

            record.index_status = 'indexed_complete'  # New final status
            record.processing_error = None  # Clear previous errors if successful
            logger.success(
                f"✅ {log_prefix}: Successfully indexed and embedded {len(chunks)} chunks for File ID {record.id}.")

        except TaskInterruptedException as tie:
            logger.warning(
                f"{log_prefix}: Embedding for File ID {record.id} was interrupted: {tie}. Will remain as 'pending_embedding'.")
            record.index_status = 'pending_embedding'  # Re-queue for next cycle
        except Exception as e_embed_update:
            logger.error(
                f"{log_prefix}: Failed to embed or update vector store for File ID {record.id}: {e_embed_update}")
            record.index_status = 'error_embedding'
            record.processing_error = f"Final embedding/Chroma update failed: {e_embed_update}"

        db_session.commit()

    def _extract_office_text(self, file_path: str) -> Optional[str]:
        """
        Extracts text content from Microsoft Office files (.docx, .pptx, .xlsx, .xls).
        """
        log_prefix = f"OfficeExtract|{os.path.basename(file_path)[:15]}"
        file_ext = os.path.splitext(file_path)[1].lower()
        extracted_content = []

        try:
            # --- Word Documents (.docx) ---
            if file_ext == '.docx':
                # Ensure 'python-docx' is installed
                import docx
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    if para.text.strip():
                        extracted_content.append(para.text)
                # You can also extract text from tables if needed
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                extracted_content.append(cell.text)

            # --- PowerPoint Presentations (.pptx) ---
            elif file_ext == '.pptx':
                # Ensure 'python-pptx' is installed
                import pptx
                prs = pptx.Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            if shape.text.strip():
                                extracted_content.append(shape.text)
                        if hasattr(shape, "has_table") and shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text.strip():
                                        extracted_content.append(cell.text)

            # --- Excel Spreadsheets (.xlsx, .xls) ---
            elif file_ext in ['.xlsx', '.xls']:
                # Ensure 'pandas' and its engines ('openpyxl', 'xlrd') are installed
                import pandas as pd
                # Read all sheets into a dictionary of DataFrames
                xls = pd.ExcelFile(file_path)
                for sheet_name in xls.sheet_names:
                    df = pd.read_excel(xls, sheet_name=sheet_name, header=None)
                    # Convert all data in the sheet to string and join
                    sheet_text = ' '.join(df.applymap(lambda x: str(x) if pd.notna(x) else '').values.flatten())
                    if sheet_text.strip():
                        extracted_content.append(f"--- Sheet: {sheet_name} ---\n{sheet_text}")

            else:
                log_worker("WARNING", f"{log_prefix}: Unsupported Office file type for text extraction: {file_ext}")
                return None

            full_text = "\n\n".join(extracted_content)
            logger.info(
                f"{log_prefix}: Successfully extracted ~{len(full_text)} characters from {os.path.basename(file_path)}.")
            return full_text

        except Exception as e:
            logger.error(f"{log_prefix}: Failed to extract text from Office file '{file_path}': {e}")
            return f"[Error extracting text from Office file: {e}]"

    async def _process_file_phase1(self, file_path: str, db_session: Session):
        """
        Phase 1: Gets metadata, hash, extracts text, AND performs an
        initial text-only embedding to make the file searchable immediately.
        This version is fully self-contained and handles errors gracefully.
        """
        if self.stop_event.is_set():
            return
        log_prefix = f"P1-File|{os.path.basename(file_path)[:20]}"

        try:
            # Step 1: Initial file validation and metadata gathering
            # =======================================================
            try:
                if not os.path.exists(file_path) or not os.path.isfile(file_path):
                    logger.debug(f"{log_prefix}: Path does not exist or is not a file. Skipping.")
                    return

                file_metadata = self._get_file_metadata(file_path)
                file_size = file_metadata.get('size_bytes')

                if file_size is None:
                    raise PermissionError("Could not retrieve file metadata, likely a permission issue.")
                if file_size > FILE_INDEX_MAX_SIZE_MB * 1024 * 1024:
                    return
                if file_size < FILE_INDEX_MIN_SIZE_KB * 1024:
                    return

            except (OSError, PermissionError) as e_stat:
                logger.warning(f"{log_prefix}: Stat failed or permission denied: {e_stat}. Logging error and skipping file.")
                # Log a specific error status to the DB and stop processing this file
                existing_rec = db_session.query(FileIndex).filter(FileIndex.file_path == file_path).first()
                if not existing_rec:
                    err_record = FileIndex(
                        file_path=file_path, file_name=os.path.basename(file_path),
                        index_status=IndexStatusEnum.error_permission,
                        processing_error=f"Initial file access/stat failed: {e_stat}"
                    )
                    db_session.add(err_record)
                    db_session.commit()
                return

            if self._wait_if_server_busy():
                return

            # Step 2: Check for modifications using MD5 hash
            # =================================================
            current_md5 = self._calculate_md5(file_path, file_size)
            if not current_md5 and file_size <= MAX_HASH_FILE_SIZE_BYTES:
                logger.warning(f"{log_prefix}: MD5 calculation failed for a hashable file size. Skipping.")
                return

            existing_record: Optional[FileIndex] = db_session.query(FileIndex).filter(
                FileIndex.file_path == file_path).first()
            values_to_update: Dict[str, Any] = {'md5_hash': current_md5}
            values_to_update.update(file_metadata)

            needs_processing = False
            if not existing_record:
                needs_processing = True
                logger.info(f"-> {log_prefix}: New file detected.")
            else:
                mtime_os_check = file_metadata.get('last_modified_os')
                hashes_match = (current_md5 is not None and existing_record.md5_hash == current_md5)
                large_file_is_unchanged = (
                    current_md5 is None and existing_record.md5_hash is None and
                    existing_record.size_bytes == file_size and
                    mtime_os_check and existing_record.last_modified_os and
                    mtime_os_check <= existing_record.last_modified_os
                )
                if hashes_match or large_file_is_unchanged:
                    if not existing_record.indexed_content and existing_record.index_status not in [IndexStatusEnum.error_read, IndexStatusEnum.error_permission, IndexStatusEnum.skipped_type]:
                        needs_processing = True
                        logger.info(f"-> {log_prefix}: Unchanged file is missing text content. Re-processing.")
                    else:
                        logger.trace(f"{log_prefix}: Unchanged and already processed. Skipping.")
                        return
                else:
                    needs_processing = True
                    logger.info(f"-> {log_prefix}: File modified. Re-processing.")

            # Step 3: Extract Text and Set Initial Status
            # ===============================================
            if needs_processing:
                content: Optional[str] = None
                file_ext = os.path.splitext(file_path)[1].lower()

                try:
                    if file_ext in OCR_TARGET_EXTENSIONS:
                        content = self._extract_text_with_ocr_fallback(file_path, file_ext)
                        if file_ext in VLM_TARGET_EXTENSIONS:
                            values_to_update['vlm_processing_status'] = 'pending_vlm'
                        values_to_update['index_status'] = IndexStatusEnum.pending_embedding
                    elif file_ext in TEXT_EXTENSIONS:
                        content = self._extract_text(file_path, file_size)
                        values_to_update['index_status'] = IndexStatusEnum.pending_embedding
                    elif file_ext in OFFICE_EXTENSIONS:
                        content = self._extract_office_text(file_path)
                        if file_ext in VLM_TARGET_EXTENSIONS:
                            values_to_update['vlm_processing_status'] = 'pending_vlm'
                        values_to_update['index_status'] = IndexStatusEnum.pending_embedding
                    else:
                        values_to_update['index_status'] = IndexStatusEnum.skipped_type

                except PermissionError:
                    logger.warning(f"{log_prefix} Permission denied during text extraction.")
                    values_to_update['index_status'] = IndexStatusEnum.error_permission
                    values_to_update['processing_error'] = "Permission denied during text extraction."
                except Exception as e_extract:
                    logger.error(f"{log_prefix} Error during text extraction: {e_extract}", exc_info=True)
                    values_to_update['index_status'] = IndexStatusEnum.error_read
                    values_to_update['processing_error'] = f"Text extraction failed: {str(e_extract)[:255]}"

                values_to_update['last_indexed_db'] = datetime.datetime.now(datetime.timezone.utc)
                
                # Step 4: Create/Update DB Record Before Embedding
                # ====================================================
                record_to_embed: Optional[FileIndex] = None
                if existing_record:
                    if values_to_update.get('vlm_processing_status') == 'pending_vlm':
                        values_to_update['latex_representation'] = None
                        values_to_update['latex_explanation'] = None
                    
                    values_to_update['indexed_content'] = content # ADDED THIS LINE
                    db_session.execute(update(FileIndex).where(FileIndex.id == existing_record.id).values(**values_to_update))
                    db_session.commit()
                    record_to_embed = existing_record
                else:
                    new_record_data = {
                        "file_path": file_path, 
                        "file_name": os.path.basename(file_path), 
                        "indexed_content": content, # ADDED THIS LINE
                        **values_to_update
                    }
                    record_to_embed = FileIndex(**new_record_data)
                    db_session.add(record_to_embed)
                    db_session.commit()
                    db_session.refresh(record_to_embed)

                # Step 5: Perform Initial Embedding
                # =====================================
                if content and record_to_embed:
                    logger.info(f"{log_prefix}: Text extracted. Proceeding to initial embedding...")
                    await self._embed_and_store(record_to_embed, content, db_session, is_final_embedding=False)
                elif not content:
                    logger.debug(f"{log_prefix}: No content extracted, skipping initial embedding.")

        except Exception as e_phase1:
            logger.error(f"{log_prefix}: Unhandled error in Phase 1 for {file_path}: {e_phase1}", exc_info=True)

    def _cleanup_llm_output(self, text: str, replacement_char: str = ' ') -> str:
        """
        Cleans up raw text output from the LLM, removing special tokens
        and extraneous whitespace.
        """
        if not isinstance(text, str):
            return ""

        # This cleanup logic should match the one in app.py's CortexThoughts class for consistency
        sanitized = text.strip()

        # Use config variables if they are available, otherwise use defaults
        start_token = CHATML_START_TOKEN if 'CHATML_START_TOKEN' in globals() else "<|im_start|>"
        end_token = CHATML_END_TOKEN if 'CHATML_END_TOKEN' in globals() else "<|im_end|>"

        # Remove start/end tokens and any "assistant\n" preamble
        if sanitized.startswith(start_token):
            sanitized = sanitized[len(start_token):].lstrip()
        if sanitized.endswith(end_token):
            sanitized = sanitized[:-len(end_token)].rstrip()

        # A regex to remove "assistant" only if it's at the very beginning of the string,
        # potentially with some whitespace.
        sanitized = re.sub(r"^\s*assistant\s*\n?", "", sanitized, flags=re.IGNORECASE)

        # Replace multiple spaces/newlines with a single character
        # The 'r' prefix makes it a raw string, preventing SyntaxWarning for '\s'
        sanitized = re.sub(rf'[{re.escape(replacement_char)}\s]+', replacement_char, sanitized)

        return sanitized.strip()

    async def _get_initial_vlm_description_async(self, image: Image.Image, session_id_for_log: str) -> Tuple[
        Optional[str], Optional[str]]:
        """
        Calls the general 'vlm' model to get an initial textual description of an image.
        Returns (description, error_message_if_any).
        """
        # Check for stop events before processing
        if self._wait_if_server_busy(): return None, "[VLM Description Skipped - Server Busy]"
        if self.stop_event.is_set(): return None, "[VLM Description Skipped - Stop Requested]"

        if not self.vlm_model:
            log_worker("WARNING", "VLM model ('vlm' role) not available for initial description.")
            return None, "[VLM Model Unavailable]"

        log_prefix = f"VLMDesc|ELP0|{session_id_for_log[:8]}"
        logger.trace(f"{log_prefix}: Getting initial VLM description for image...")

        description_output: Optional[str] = None
        error_output: Optional[str] = None

        try:
            # 1. Convert PIL Image to base64 data URI for the VLM prompt
            buffered = BytesIO()
            image.save(buffered, format="PNG")  # PNG is a good lossless format for this
            img_b64_str = base64.b64encode(buffered.getvalue()).decode('utf-8')
            image_uri = f"data:image/png;base64,{img_b64_str}"

            # 2. Prepare the prompt and messages for the VLM
            # PROMPT_VLM_INITIAL_ANALYSIS should be defined in config.py
            # e.g., PROMPT_VLM_INITIAL_ANALYSIS = "Describe this image, focusing on text, formulas, or diagrams."
            image_content = {"type": "image_url", "image_url": {"url": image_uri}}
            text_content = {"type": "text", "text": PROMPT_VLM_INITIAL_ANALYSIS}
            vlm_messages = [HumanMessage(content=[image_content, text_content])]

            # 3. Create the Langchain chain
            vlm_chain = self.vlm_model | StrOutputParser()

            # 4. Call the LLM (VLM)
            # _call_llm_with_timing is synchronous, so wrap it for our async context
            timing_data = {"session_id": session_id_for_log, "mode": "vlm_initial_description", "execution_time_ms": 0}

            response_text = await asyncio.to_thread(
                self.provider._call_llm_with_timing,  
                vlm_chain,
                vlm_messages,  # Pass messages directly as input to the model in the chain
                timing_data,
                priority=ELP0  # This is a background task
            )

            if response_text and not (
                    isinstance(response_text, str) and "ERROR" in response_text and "Traceback" in response_text):
                description_output = self._cleanup_llm_output(response_text.strip())
                logger.trace(f"{log_prefix}: VLM description successful. Snippet: '{description_output[:100]}...'")
            else:
                error_output = f"[VLM description call failed or returned error: {response_text}]"
                logger.warning(f"{log_prefix} {error_output}")

        except TaskInterruptedException as tie:
            logger.warning(f"🚦 {log_prefix} VLM description INTERRUPTED: {tie}")
            error_output = "[VLM Description Interrupted]"
        except Exception as e:
            logger.error(f"{log_prefix} VLM description call failed: {e}", exc_info=True)
            error_output = f"[VLM Description Error: {str(e)[:100]}]"

        return description_output, error_output

    async def _process_pending_vlm_files(self, db_session: Session):
        """
        Phase 2: Processes files marked as 'pending_vlm' or 'pending_conversion'.
        1. Converts files to images if necessary (e.g., Office docs).
        2. Runs general VLM for description and LaTeX-OCR for math on each page image.
        3. Saves this new text to the database.
        4. Calls the final embedding helper to create and store vectors for the combined text.
        """
        if not self.vlm_model and not self.latex_ocr_model:
            logger.warning("Skipping Phase 2 processing: VLM and/or LatexOCR models are not available.")
            return

        logger.info("🔬 Starting Phase 2 VLM/OCR Cycle...")

        # Process a limited number of files per cycle to avoid getting stuck on a long task
        # and to allow the main scan loop to run periodically.
        processed_in_cycle = 0
        max_to_process_in_cycle = 50

        while not self.stop_event.is_set() and processed_in_cycle < max_to_process_in_cycle:
            if self._wait_if_server_busy():
                logger.info("Phase 2: Pausing due to busy server.")
                break

            pending_vlm_files: List[FileIndex] = []
            try:
                # Query for files needing VLM processing
                pending_vlm_files = db_session.query(FileIndex).filter(
                    FileIndex.vlm_processing_status.in_(['pending_vlm', 'pending_conversion'])
                ).limit(10).all()  # Process in batches of 10
            except Exception as e_query:
                logger.error(f"Phase 2: Error querying for pending VLM files: {e_query}")
                break

            if not pending_vlm_files:
                logger.info("Phase 2: No more files pending VLM/OCR processing in this batch.")
                break  # Exit cycle if no more files are found

            record: FileIndex
            for record in pending_vlm_files:
                if self.stop_event.is_set():
                    logger.info("Phase 2: Stop signal received during batch processing.")
                    break

                processed_in_cycle += 1
                log_prefix = f"P2-VLM|{os.path.basename(record.file_path)[:15]}" #type: ignore
                logger.info(f"--> {log_prefix}: Starting Phase 2 processing for File ID {record.id}")

                images: Optional[List[Image.Image]] = None
                temp_pdf_to_process: Optional[str] = None

                try:
                    # --- FIX APPLIED HERE ---
                    # Get the file extension from the file_path attribute
                    file_ext = os.path.splitext(record.file_path)[1].lower() #type: ignore

                    # Step 1: Convert file to a list of PIL Images
                    if file_ext == '.pdf':
                        images = await asyncio.to_thread(self._convert_pdf_to_images, record.file_path) #type: ignore
                    elif file_ext in OFFICE_EXTENSIONS:
                        temp_pdf_to_process = await asyncio.to_thread(self._convert_office_to_pdf, record.file_path) #type: ignore
                        if temp_pdf_to_process and os.path.exists(temp_pdf_to_process):
                            images = await asyncio.to_thread(self._convert_pdf_to_images, temp_pdf_to_process)
                    elif file_ext in VLM_TARGET_EXTENSIONS:  # For single images like PNG, JPG
                        images = [Image.open(record.file_path)] #type: ignore

                    if not images:
                        raise ValueError("Failed to convert or load file into images for VLM/OCR processing.")

                    page_vlm_descriptions: List[str] = []
                    page_latex_ocr_outputs: List[str] = []
                    page_errors: List[str] = []

                    # Step 2: Process each page/image with VLM and LaTeX-OCR
                    for page_num, page_image in enumerate(images):
                        session_id_for_page = f"p2_vlm_{record.id}_p{page_num}"

                        # Get general description from main VLM model
                        if self.vlm_model:
                            vlm_description, vlm_err = await self._get_initial_vlm_description_async(page_image,
                                                                                                     session_id_for_page)
                            if vlm_err: page_errors.append(f"Page {page_num + 1} VLM Error: {vlm_err}")
                            if vlm_description: page_vlm_descriptions.append(vlm_description)

                        # Get LaTeX math from specialized LatexOCR model
                        if self.latex_ocr_model:
                            latex_math = await asyncio.to_thread(self._get_latex_from_image, page_image)
                            if latex_math:
                                if "[LatexOCR Error:" not in latex_math:
                                    page_latex_ocr_outputs.append(latex_math)
                                else:
                                    page_errors.append(f"Page {page_num + 1} LatexOCR Error: {latex_math}")

                    # Step 3: Combine results and update the record in the database session
                    final_vlm_desc = "\n\n--- Page Break ---\n\n".join(
                        page_vlm_descriptions).strip() if page_vlm_descriptions else None
                    final_latex_math = "\n\n".join(page_latex_ocr_outputs).strip() if page_latex_ocr_outputs else None
                    final_errors = "\n".join(page_errors) if page_errors else None

                    record.latex_representation = final_latex_math  # Column for specialized math OCR
                    record.latex_explanation = final_vlm_desc  # Column for general VLM description
                    record.processing_error = final_errors

                    if final_vlm_desc or final_latex_math:
                        record.vlm_processing_status = 'success_vlm_ocr' if not page_errors else 'partial_vlm_ocr_error'
                    else:
                        record.vlm_processing_status = 'error_vlm_ocr'

                    logger.info(
                        f"{log_prefix}: VLM/OCR text extraction complete. Status: {record.vlm_processing_status}. Proceeding to final embedding.")

                    # Step 4: Call the unified embedding function now that all text is available
                    # This function will handle chunking, embedding, updating Chroma, and setting final status.
                    await self._embed_and_update_vector_store(record, db_session)
                    # The commit is handled inside _embed_and_update_vector_store

                except Exception as e_vlm_process:
                    error_message = f"Failed during Phase 2 processing loop: {str(e_vlm_process)}"
                    logger.error(f"{log_prefix}: {error_message}", exc_info=True)
                    # Update DB with error status for this record
                    db_session.execute(update(FileIndex).where(FileIndex.id == record.id).values(
                        vlm_processing_status='error_vlm_ocr',
                        processing_error=error_message[:1000],
                        last_indexed_db=datetime.datetime.now(datetime.timezone.utc)
                    ))
                    db_session.commit()
                finally:
                    # Clean up temporary PDF from Office conversion if it exists
                    if temp_pdf_to_process and os.path.exists(temp_pdf_to_process):
                        try:
                            os.remove(temp_pdf_to_process)
                        except Exception as e_rm_tmp:
                            logger.warning(
                                f"{log_prefix}: Could not remove temp PDF '{temp_pdf_to_process}': {e_rm_tmp}")

            if self.stop_event.is_set(): break  # Exit outer while loop if stopped during batch

        logger.info("Phase 2 VLM/OCR Cycle Finished.")

    async def _process_pending_embeddings(self, db_session: Session):
        """
        Phase 3: Processes files marked as 'pending_embedding'.
        These are files that have had their text extracted and are ready for the
        final combined embedding and vector store update. This typically includes
        plain text files or OCR'd images that did not require VLM analysis.
        """
        logger.info("🔬 Starting Phase 3 Final Embedding Cycle...")
        processed_in_cycle = 0
        max_to_process_in_cycle = 100  # Process up to 100 files per cycle

        while not self.stop_event.is_set() and processed_in_cycle < max_to_process_in_cycle:
            if self._wait_if_server_busy():
                logger.info("Phase 3: Pausing due to busy server.")
                break

            pending_embedding_files: List[FileIndex] = []
            try:
                # Query for files that are ready for the final embedding step
                pending_embedding_files = db_session.query(FileIndex).filter(
                    FileIndex.index_status == 'pending_embedding'
                ).limit(20).all()  # Process in batches of 20
            except Exception as e_query_embed:
                logger.error(f"Phase 3: Error querying for pending embedding files: {e_query_embed}")
                break

            if not pending_embedding_files:
                logger.info("Phase 3: No more files pending final embedding in this batch.")
                break

            record: FileIndex
            for record in pending_embedding_files:
                if self.stop_event.is_set():
                    logger.info("Phase 3: Stop signal received during batch processing.")
                    break

                processed_in_cycle += 1
                log_prefix = f"P3-Embed|{os.path.basename(record.file_path)[:15]}" #type: ignore
                logger.info(f"--> {log_prefix}: Starting final embedding for File ID {record.id}")

                try:
                    # Call the unified embedding helper function. This function handles
                    # combining text, chunking, embedding, and updating ChromaDB.
                    await self._embed_and_update_vector_store(record, db_session)
                except Exception as e_final_embed:
                    logger.error(
                        f"{log_prefix}: Unhandled error during final embedding process for File ID {record.id}: {e_final_embed}",
                        exc_info=True)
                    record.index_status = 'error_embedding'
                    record.processing_error = f"Final embedding orchestration failed: {str(e_final_embed)[:255]}"
                    db_session.commit()

            if self.stop_event.is_set(): break

        logger.info("Phase 3 Final Embedding Cycle Finished.")

    # --- (Existing run method - orchestrates Phase 1 then Phase 2) ---
    def run(self):
        """
        The main loop for the file indexer thread. This version uses separate DB sessions
        for each phase to increase resilience.
        """

        if not DATABASE_AVAILABLE:
            logger.critical(
                f"🛑 {self.thread_name} cannot start because database components are missing. Thread is shutting down.")
            return  # Exit the run method immediately

        logger.info(f"✅ {self.thread_name} started (Multi-Phase, Isolated Session Logic).")
        
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)

        try:
            loop.run_until_complete(initialize_global_file_index_vectorstore(self.provider))
        except Exception as e_init_vs:
            logger.error(f"CRITICAL: Vector store loading failed: {e_init_vs}. Thread will exit.")
            return

        while not self.stop_event.is_set():
            cycle_start_time = time.monotonic()

            # --- Phase 1 ---
            try:
                with SessionLocal() as db_session: # Session for Phase 1
                    logger.info(f"--- {self.thread_name}: Starting Phase 1 (Scan & Text Extraction) ---")
                    for root_path in self._get_root_paths():
                        #self._scan_directory(root_path, db_session)
                        self._scan_directory(root_path, db_session, loop)
                        if self.stop_event.is_set(): break
            except Exception as e_phase1:
                logger.error(f"💥 Unhandled error in Phase 1: {e_phase1}", exc_info=True)
            
            if self.stop_event.is_set(): break

            # --- Phase 2 ---
            try:
                with SessionLocal() as db_session: # New Session for Phase 2
                    logger.info(f"--- {self.thread_name}: Starting Phase 2 (VLM & LaTeX-OCR) ---")
                    loop.run_until_complete(self._process_pending_vlm_files(db_session))
            except Exception as e_phase2:
                logger.error(f"💥 Unhandled error in Phase 2: {e_phase2}", exc_info=True)

            if self.stop_event.is_set(): break

            # --- Phase 3 ---
            try:
                with SessionLocal() as db_session: # New Session for Phase 3
                    logger.info(f"--- {self.thread_name}: Starting Phase 3 (Final Embedding) ---")
                    loop.run_until_complete(self._process_pending_embeddings(db_session))
            except Exception as e_phase3:
                logger.error(f"💥 Unhandled error in Phase 3: {e_phase3}", exc_info=True)

            cycle_duration = time.monotonic() - cycle_start_time
            logger.info(f"--- {self.thread_name}: Full pipeline cycle finished in {cycle_duration:.2f} seconds. ---")
            
            wait_time = FILE_INDEXER_IDLE_WAIT_SECONDS
            logger.debug(f"{self.thread_name} waiting for {wait_time}s before next full cycle...")
            if self.stop_event.wait(timeout=wait_time):
                break

        try:
            loop.close()
        except Exception as e_loop_close:
            logger.warning(f"Error closing asyncio loop in {self.thread_name}: {e_loop_close}")

        logger.info(f"🛑 {self.thread_name} has been shut down.")


def _locked_initialization_task(provider_ref: CortexEngine) -> Dict[str, Any]:
    global global_file_index_vectorstore

    # These globals() calls assume these variables are available from CortexConfiguration.py
    _persist_dir_to_use = globals().get("FILE_INDEX_CHROMA_PERSIST_DIR")
    _collection_name_to_use = globals().get("FILE_INDEX_CHROMA_COLLECTION_NAME")

    # Allow one retry attempt if we need to clear the old DB due to a settings conflict.
    for attempt in range(2):
        task_status: str = "unknown_error_before_processing"
        task_message: str = "Initialization did not complete as expected."
        initialization_succeeded_or_known_empty: bool = False
        db_session: Optional[Session] = None

        overall_start_time: float = time.monotonic()
        logger.info(
            f">>> FileIndex VS Init (Attempt {attempt + 1}/2): Attempting to acquire _file_index_vs_init_lock... <<<")

        try:
            with _file_index_vs_init_lock:
                lock_acquired_time: float = time.monotonic()
                logger.info(
                    f">>> FileIndex VS Init: ACQUIRED _file_index_vs_init_lock (waited {lock_acquired_time - overall_start_time:.3f}s). <<<")

                if _file_index_vs_initialized_event.is_set():
                    if global_file_index_vectorstore is not None:
                        logger.info(
                            ">>> FileIndex VS Init: SKIPPING - _file_index_vs_initialized_event is ALREADY SET and VS exists. <<<")
                        return {"status": "skipped_already_initialized",
                                "message": "Initialization event was already set and VS exists."}
                    else:
                        logger.warning(
                            ">>> FileIndex VS Init: Event was set, but global_file_index_vectorstore is None. This indicates a previous partial/failed init. Clearing event and re-attempting full initialization under lock.")
                        _file_index_vs_initialized_event.clear()

                # STAGE 0: Validate Provider and Embeddings
                logger.info(">>> FileIndex VS Init: Stage 0: Validating Provider and Embeddings. <<<")
                if not provider_ref or not provider_ref.embeddings:
                    task_status = "error_provider_missing"
                    task_message = "CortexEngine or its embeddings module is missing. Cannot initialize FileIndex VS."
                    logger.error(task_message)
                    initialization_succeeded_or_known_empty = True
                    global_file_index_vectorstore = None
                    if not _file_index_vs_initialized_event.is_set(): _file_index_vs_initialized_event.set()
                    return {"status": task_status, "message": task_message}

                if not _persist_dir_to_use or not _collection_name_to_use:
                    task_status = "error_config_missing_chroma_paths"
                    task_message = "Chroma DB persist directory or collection name not configured."
                    logger.error(task_message)
                    initialization_succeeded_or_known_empty = True
                    global_file_index_vectorstore = None
                    if not _file_index_vs_initialized_event.is_set(): _file_index_vs_initialized_event.set()
                    return {"status": task_status, "message": task_message}

                # Define consistent client settings to prevent the error in the first place.
                chroma_settings = Settings(anonymized_telemetry=False)

                # STAGE 1: Attempt to load persisted Chroma DB
                logger.info(
                    f">>> FileIndex VS Init: Stage 1: Attempting to load persisted Chroma VS from: '{_persist_dir_to_use}' collection '{_collection_name_to_use}' <<<")

                should_rebuild_from_sql = True
                if os.path.exists(_persist_dir_to_use) and os.path.isdir(_persist_dir_to_use):
                    # Check for actual Chroma data files before attempting to load
                    if any(fname.endswith(('.sqlite3', '.duckdb', '.parquet')) for fname in
                           os.listdir(_persist_dir_to_use)):
                        logger.info(f"Found persisted Chroma data files. Attempting to load collection...")
                        loaded_store = Chroma(
                            collection_name=_collection_name_to_use,
                            embedding_function=provider_ref.embeddings,
                            persist_directory=_persist_dir_to_use,
                            client_settings=chroma_settings
                        )
                        collection_count = loaded_store._collection.count()

                        if collection_count > 0:
                            global_file_index_vectorstore = loaded_store
                            task_status = "success_loaded_from_persist"
                            task_message = f"Successfully loaded FileIndex VS from '{_persist_dir_to_use}' with {collection_count} items."
                            should_rebuild_from_sql = False
                            logger.success(task_message)
                        else:
                            logger.warning(
                                f"Loaded persisted Chroma from '{_persist_dir_to_use}', but it's EMPTY. Will attempt rebuild from SQL.")
                    else:
                        logger.info(f"No Chroma data files found in '{_persist_dir_to_use}'. Will build from SQL DB.")
                else:
                    logger.info(f"Persist directory '{_persist_dir_to_use}' does not exist. Will build from SQL DB.")
                    os.makedirs(_persist_dir_to_use, exist_ok=True)

                if should_rebuild_from_sql:
                    logger.info(">>> FileIndex VS Init: Proceeding to rebuild from SQL Database. <<<")
                    db_session = SessionLocal()
                    if not db_session:
                        raise RuntimeError("Failed to create DB session for rebuild phase.")

                    logger.info(f">>> FileIndex VS Init: Stage 2: Querying SQL DB for records with embeddings... <<<")
                    # This query matches the logic in your original code
                    indexed_files = db_session.query(FileIndex).filter(
                        FileIndex.embedding_json.isnot(None),
                        FileIndex.indexed_content.isnot(None),
                        FileIndex.index_status.in_(['indexed_text', 'success', 'partial_vlm_error'])
                    ).all()

                    total_records_from_db = len(indexed_files)
                    logger.info(
                        f"Stage 2 (DB Query) completed. Found {total_records_from_db} candidates with embeddings in SQL.")

                    if not indexed_files:
                        logger.warning("No files with embeddings in SQL DB to rebuild. Creating empty persistent VS.")
                        global_file_index_vectorstore = Chroma(
                            collection_name=_collection_name_to_use,
                            embedding_function=provider_ref.embeddings,
                            persist_directory=_persist_dir_to_use,
                            client_settings=chroma_settings
                        )
                        task_status = "success_empty_db_rebuild"
                        task_message = "Init (rebuild): No files with embeddings in SQL. Created empty VS."
                    else:
                        logger.info(
                            f">>> FileIndex VS Init: Stage 3: Processing {total_records_from_db} DB records for Chroma... <<<")
                        texts_for_vs, embeddings_for_vs, metadatas_for_vs, ids_for_vs = [], [], [], []
                        processed_records_for_chroma = 0

                        for db_record in indexed_files:
                            try:
                                rec_id_val = getattr(db_record, 'id')
                                emb_json_str = getattr(db_record, 'embedding_json', None)
                                content_str = getattr(db_record, 'indexed_content', None)

                                if not emb_json_str or not content_str or not content_str.strip():
                                    continue

                                vector = json.loads(emb_json_str)
                                if not (isinstance(vector, list) and vector and all(
                                        isinstance(val, (float, int)) for val in vector)):
                                    continue

                                texts_for_vs.append(content_str)
                                embeddings_for_vs.append([float(val) for val in vector])

                                # Replicating metadata from original code
                                metadatas_for_vs.append({
                                    "source": getattr(db_record, 'file_path', "UnknownPath"),
                                    "file_id": int(rec_id_val) if str(rec_id_val).isdigit() else -1,
                                    "file_name": getattr(db_record, 'file_name', "UnknownFile"),
                                    "last_modified": str(getattr(db_record, 'last_modified_os', "N/A")),
                                    "index_status": str(getattr(db_record, 'index_status', "UnknownStatus")),
                                    "mime_type": getattr(db_record, 'mime_type', "UnknownMIME")
                                })
                                ids_for_vs.append(f"file_{rec_id_val}")
                                processed_records_for_chroma += 1
                            except Exception as e_prep:
                                logger.warning(
                                    f"Skipping record ID {getattr(db_record, 'id', 'N/A')} due to data preparation error: {e_prep}")

                        logger.info(
                            f"Stage 3 (DB Record Processing) completed. Prepared {processed_records_for_chroma} valid items.")

                        if processed_records_for_chroma > 0:
                            logger.info(
                                f">>> FileIndex VS Init: Stage 4: Populating Chroma with {processed_records_for_chroma} items... <<<")

                            # Create a new store for population
                            temp_chroma_store = Chroma(
                                collection_name=_collection_name_to_use,
                                embedding_function=provider_ref.embeddings,
                                persist_directory=_persist_dir_to_use,
                                client_settings=chroma_settings
                            )

                            # Add data to the collection
                            temp_chroma_store._collection.add(
                                ids=ids_for_vs,
                                embeddings=embeddings_for_vs,
                                metadatas=metadatas_for_vs,
                                documents=texts_for_vs
                            )

                            global_file_index_vectorstore = temp_chroma_store
                            task_status = "success_rebuilt_from_sql"
                            task_message = f"FileIndex VS rebuilt from SQL with {processed_records_for_chroma} items."
                            logger.success(task_message)
                        else:
                            logger.warning("No valid data with embeddings from SQL. Creating empty persistent VS.")
                            global_file_index_vectorstore = Chroma(
                                collection_name=_collection_name_to_use,
                                embedding_function=provider_ref.embeddings,
                                persist_directory=_persist_dir_to_use,
                                client_settings=chroma_settings
                            )
                            task_status = "success_no_valid_data_rebuild"
                            task_message = "Init (rebuild): No valid data from SQL. Created empty VS."

                initialization_succeeded_or_known_empty = True
                break  # <-- IMPORTANT: Exit the retry loop on success

        except ValueError as e:
            # This is the core of the requested change
            if "An instance of Chroma already exists" in str(e) and attempt == 0:
                logger.warning(f"Chroma conflict detected: {e}. Deleting directory and retrying...")
                # The lock is already held, so this delete operation is thread-safe
                if _persist_dir_to_use and os.path.isdir(_persist_dir_to_use):
                    try:
                        shutil.rmtree(_persist_dir_to_use)
                        logger.info(f"Successfully deleted Chroma directory: {_persist_dir_to_use}")
                    except Exception as rm_e:
                        task_status = "critical_error_delete_failed"
                        task_message = f"CRITICAL: Failed to delete Chroma directory {_persist_dir_to_use}: {rm_e}"
                        logger.error(task_message)
                        # If we can't delete the directory, we must stop trying.
                        break
                # Continue to the next iteration of the loop to retry initialization.
                continue
            else:
                # If this is another ValueError or it's the second attempt, fail permanently.
                task_status = "critical_error_ValueError"
                task_message = f"CRITICAL ERROR during FileIndex VS init (ValueError): {e}"
                logger.error(task_message, exc_info=True)
                global_file_index_vectorstore = None
                initialization_succeeded_or_known_empty = False
                break  # Exit loop

        except Exception as e_init_critical:
            task_status = "critical_error_overall_init"
            task_message = f"CRITICAL ERROR during FileIndex VS init: {e_init_critical}"
            logger.error(task_message, exc_info=True)
            global_file_index_vectorstore = None
            initialization_succeeded_or_known_empty = False
            break  # Exit loop

        finally:
            if db_session:
                try:
                    db_session.close()
                except Exception as e_close:
                    logger.warning(f"Error closing DB session in FileIndex VS init: {e_close}")

            # This block runs at the end of each attempt's try/except/finally
            if attempt == 1 and not initialization_succeeded_or_known_empty:
                logger.error(">>> All initialization attempts failed. The vector store will be unavailable. <<<")

    # This block runs after the loop is fully exhausted or broken out of
    if initialization_succeeded_or_known_empty:
        if not _file_index_vs_initialized_event.is_set():
            _file_index_vs_initialized_event.set()
            logger.info(f">>> SET _file_index_vs_initialized_event. Final Status: {task_status}. <<<")
    else:
        logger.error(f">>> _file_index_vs_initialized_event NOT SET due to critical error. Status: {task_status}. <<<")

    logger.info(f">>> FileIndex VS Init: Task finished. Final Status: '{task_status}', Message: \"{task_message}\" <<<")
    return {"status": task_status, "message": task_message}


async def initialize_global_file_index_vectorstore(provider: CortexEngine):
    global global_file_index_vectorstore  # Still need this for the outer function's scope
    logger.info(">>> Entered initialize_global_file_index_vectorstore (async version) <<<")

    if not provider:
        logger.error("  INIT_VS_ERROR (Outer): CortexEngine instance is None! Cannot initialize file index vector store.")
        return
    if not provider.embeddings:
        logger.error(
            "  INIT_VS_ERROR (Outer): CortexEngine.embeddings is None! Cannot initialize file index vector store.")
        return
    logger.info(f"  INIT_VS_INFO (Outer): Using CortexEngine embeddings of type: {type(provider.embeddings)}")

    if _file_index_vs_initialized_event.is_set():
        logger.info(">>> SKIPPING FileIndex VS Init (Outer): _file_index_vs_initialized_event is ALREADY SET. <<<")
        return

    try:
        # Pass the provider instance to the threaded function _locked_initialization_task
        result_dict = await asyncio.to_thread(_locked_initialization_task, provider)

        status_from_thread = result_dict.get("status", "unknown_thread_exit_status")
        message_from_thread = result_dict.get("message", "No message from thread.")
        logger.info(
            f"Locked initialization task for FileIndex VS completed with status: '{status_from_thread}'. Message: '{message_from_thread}'")

        # The event _file_index_vs_initialized_event is set (or not set) inside _locked_initialization_task's finally block.
        # No need to set it again here unless _locked_initialization_task failed to be scheduled by to_thread.

    except Exception as e_thread_task:
        logger.error(f"Exception running _locked_initialization_task via asyncio.to_thread: {e_thread_task}")
        logger.exception("FileIndex VS _locked_initialization_task Thread Execution Traceback:")
        # If to_thread itself fails, the event won't be set by the inner task.
        # We might choose to set it here to prevent hangs, indicating init was attempted but failed at thread level.
        if not _file_index_vs_initialized_event.is_set():
            logger.warning(
                ">>> Setting _file_index_vs_initialized_event due to THREAD EXECUTION ERROR (VS is None). <<<")
            _file_index_vs_initialized_event.set()
            global_file_index_vectorstore = None  # Ensure it's None

    logger.info(
        f">>> EXITED initialize_global_file_index_vectorstore. Event set: {_file_index_vs_initialized_event.is_set()} Global VS is None: {global_file_index_vectorstore is None}<<<")


def get_global_file_index_vectorstore(timeout_seconds=60) -> Optional[Chroma]:
    if not _file_index_vs_initialized_event.is_set():
        logger.warning(f"Global file index VS not ready, waiting up to {timeout_seconds}s for initialization event...")
        event_was_set = _file_index_vs_initialized_event.wait(timeout=timeout_seconds)
        if not event_was_set:
            logger.error(f"Timeout waiting for file index VS initialization after {timeout_seconds}s.")
            return None
        logger.info("File index VS initialization event received after waiting.")

    # Event is now set (or was already set)
    if global_file_index_vectorstore is None:
        logger.warning(
            "File index VS event is set, but global_file_index_vectorstore is still None. Data might be empty or init issue.")
    return global_file_index_vectorstore


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="File Indexer Test CLI")
    parser.add_argument("--test", action="store_true", help="Run test initialization and optional search.")
    parser.add_argument("--search_query", type=str, default=None, help="Query string to search the file index vector store.")
    cli_args = parser.parse_args()

    if cli_args.test:
        logger.remove() # Remove default loguru handler
        logger.add(sys.stderr, level="DEBUG") # Add a new one with DEBUG level for testing
        logger.info("--- File Indexer Test Mode ---")

        # 1. Initialize Database (critical for SessionLocal and schema)
        try:
            logger.info("Initializing database via database.init_db()...")
            init_db() # This function from database.py handles Alembic migrations etc.
            logger.info("Database initialization complete.")
        except Exception as e_db_init:
            logger.error(f"Failed to initialize database: {e_db_init}")
            sys.exit(1)

        # 2. Initialize CortexEngine (needed for embeddings)
        #    Ensure PROVIDER is set in your environment or config.py for this to work.
        #    The CortexEngine constructor might try to load all models, which could be slow.
        #    For testing search, only provider.embeddings is strictly needed by initialize_global_file_index_vectorstore.
        test_cortex_backbone_provider: Optional[CortexEngine] = None
        try:
            logger.info(f"Initializing CortexEngine (Provider: {PROVIDER})...") # PROVIDER from CortexConfiguration.py
            test_cortex_backbone_provider = CortexEngine(PROVIDER)
            if not test_cortex_backbone_provider.embeddings:
                raise ValueError("CortexEngine initialized, but embeddings are not available.")
            logger.info("CortexEngine initialized successfully for embeddings.")
        except Exception as e_cortex_backbone_provider:
            logger.error(f"Failed to initialize CortexEngine: {e_cortex_backbone_provider}")
            sys.exit(1)

        # 3. Initialize Global File Index Vector Store (runs the async function)
        try:
            logger.info("Initializing global file index vector store (async task)...")
            asyncio.run(initialize_global_file_index_vectorstore(test_cortex_backbone_provider)) # Run the async init
            logger.info("Global file index vector store initialization process completed.")
        except Exception as e_vs_init:
            logger.error(f"Error during file index vector store initialization: {e_vs_init}")
            sys.exit(1)

        # 4. Get the vector store instance
        file_vs = get_global_file_index_vectorstore(timeout_seconds=10) # Wait a bit if still initializing
        if file_vs:
            logger.success("Successfully retrieved global file index vector store instance.")
            try:
                collection_count = file_vs._collection.count() 
                logger.info(f"File index Chroma collection contains {collection_count} items.")
            except Exception as e_count:
                logger.warning(f"Could not get count from Chroma collection: {e_count}")

            # 5. Perform Search if query provided
            if cli_args.search_query:
                query = cli_args.search_query
                logger.info(f"Performing search with query: '{query}'")
                try:
                    # Chroma's similarity_search takes the query string directly
                    # It will use the embedding_function configured during Chroma init.
                    search_results = file_vs.similarity_search_with_relevance_scores(query, k=5) # Get top 5 with scores
                    if search_results:
                        logger.info(f"Found {len(search_results)} results:")
                        for i, (doc, score) in enumerate(search_results):
                            print(f"\n--- Result {i+1} (Score: {score:.4f}) ---")
                            print(f"  Source: {doc.metadata.get('source', 'N/A')}")
                            print(f"  File Name: {doc.metadata.get('file_name', 'N/A')}")
                            print(f"  Content Snippet:\n    {doc.page_content[:300].replace(os.linesep, ' ')}...")
                    else:
                        logger.info("No results found for the query.")
                except Exception as e_search:
                    logger.error(f"Error during vector search: {e_search}")
            else:
                logger.info("No search query provided. Skipping search.")
        else:
            logger.error("Failed to get global file index vector store instance after initialization attempt.")

        logger.info("--- File Indexer Test Mode Finished ---")
    else:
        print("Run with --test to initialize and optionally --search_query \"your query\" to search.")

